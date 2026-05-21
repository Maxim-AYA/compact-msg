"""Пилот: один книжный (A4 portrait) слайд с таблицей «Документы подрядчиков» из Битрикс-Диска.

Запуск:
    python build_docs_portrait_preview.py --config <path-to-config.json>

На выходе: PPTX + PDF рядом друг с другом во временной папке, для визуальной оценки.
"""
import argparse, json, os, sys, tempfile, time
from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

A4_W_EMU = Cm(21).emu
A4_H_EMU = Cm(29.7).emu


def _read_webhook():
    path = r"C:\Авраменко\Claude Code Projects\.bitrix-webhook"
    with open(path, encoding="utf-8") as f:
        for line in f:
            s = line.strip()
            if not s or s.startswith("#"):
                continue
            wh = s.split()[0]
            break
    if not wh.startswith("http"):
        wh = f"https://kspb.bitrix24.ru/rest/2958/{wh}/"
    if not wh.endswith("/"):
        wh += "/"
    return wh


def bitrix_download_disk_file(file_id, out_path):
    import urllib.request, json as _json
    webhook = _read_webhook()
    info_url = f"{webhook}disk.file.get?id={file_id}"
    with urllib.request.urlopen(info_url, timeout=30) as r:
        data = _json.loads(r.read().decode("utf-8"))
    dl = data["result"]["DOWNLOAD_URL"]
    urllib.request.urlretrieve(dl, out_path)
    return out_path


def screenshot_range(xlsx_path, sheet_name, range_addr, out_png):
    import pythoncom, win32com.client
    from PIL import ImageGrab
    pythoncom.CoInitialize()
    excel = win32com.client.DispatchEx("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False
    try:
        wb = excel.Workbooks.Open(xlsx_path, ReadOnly=True)
        ws = wb.Sheets(sheet_name)
        ws.Activate()
        try: excel.ActiveWindow.View = 1
        except Exception: pass
        rng = ws.Range(range_addr)
        rng.CopyPicture(Appearance=1, Format=2)
        time.sleep(0.6)
        img = ImageGrab.grabclipboard()
        if img is None:
            raise RuntimeError("Excel CopyPicture produced no clipboard image")
        img.save(out_png, "PNG")
        wb.Close(SaveChanges=False)
    finally:
        excel.Quit()
        pythoncom.CoUninitialize()


def convert_pptx_to_pdf(pptx, pdf):
    import pythoncom, win32com.client
    pythoncom.CoInitialize()
    pp = win32com.client.DispatchEx("PowerPoint.Application")
    try:
        d = pp.Presentations.Open(pptx, WithWindow=False)
        d.SaveAs(pdf, 32)  # ppSaveAsPDF
        d.Close()
    finally:
        pp.Quit()
        pythoncom.CoUninitialize()


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--config", required=True)
    args = ap.parse_args()

    with open(args.config, encoding="utf-8") as f:
        cfg = json.load(f)

    file_id = cfg.get("docs_podryadchikov_file_id")
    if not file_id:
        sys.exit("config: docs_podryadchikov_file_id не задан")

    project_short = cfg.get("project_short", "")
    project_key = cfg.get("name", "проект").replace(" ", "_")

    tmp = tempfile.gettempdir()
    xlsx = os.path.join(tmp, f"docs_portrait_{project_key}.xlsx")
    bitrix_download_disk_file(file_id, xlsx)

    import openpyxl
    from openpyxl.utils import get_column_letter
    wb = openpyxl.load_workbook(xlsx, data_only=True)
    ws = wb.active
    sheet = ws.title
    rng = f"A1:{get_column_letter(ws.max_column)}{ws.max_row}"
    wb.close()

    png = os.path.join(tmp, f"docs_portrait_{project_key}.png")
    screenshot_range(xlsx, sheet, rng, png)

    # build A4 portrait PPTX with one slide
    prs = Presentation()
    prs.slide_width = A4_W_EMU
    prs.slide_height = A4_H_EMU
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # title
    title_h = Emu(700000)
    title_box = slide.shapes.add_textbox(Emu(300000), Emu(200000),
                                         A4_W_EMU - Emu(600000), title_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"ДОКУМЕНТЫ ПОДРЯДЧИКОВ ({project_short})"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # picture
    from PIL import Image
    with Image.open(png) as im:
        iw, ih = im.size
    ml = Emu(200000)
    aw = A4_W_EMU - ml * 2
    top = Emu(1000000)
    ah = A4_H_EMU - top - Emu(300000)
    scale = min(aw / iw, ah / ih)
    nw = int(iw * scale); nh = int(ih * scale)
    pl = (A4_W_EMU - nw) // 2
    pt = top + (ah - nh) // 2
    slide.shapes.add_picture(png, pl, pt, width=nw, height=nh)

    out_dir = tmp
    pptx_out = os.path.join(out_dir, f"docs_portrait_preview_{project_key}.pptx")
    pdf_out = os.path.join(out_dir, f"docs_portrait_preview_{project_key}.pdf")
    prs.save(pptx_out)
    print(f"PPTX: {pptx_out}")
    convert_pptx_to_pdf(pptx_out, pdf_out)
    print(f"PDF:  {pdf_out}")


if __name__ == "__main__":
    main()
