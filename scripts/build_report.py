"""Universal МСГ critical-lag report builder.

Reads project parameters from a JSON config and pre-extracted items.json,
produces .pptx (and prints next-step paths). PDF conversion is a separate
step (scripts/convert_pdf.py).

Usage:
    python build_report.py \\
        --config <project>/config.json \\
        --items <items.json> \\
        --week 19 \\
        --period "04.05–10.05.2026" \\
        --today 2026-05-08 \\
        [--out-pptx <path>] [--out-pdf-name <basename>]
"""
import argparse, sys, os, json
import datetime as dt
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---------------------------- args & config ----------------------------
ap = argparse.ArgumentParser()
ap.add_argument("--config", required=True, help="Project config JSON")
ap.add_argument("--items",  required=True, help="Pre-extracted items JSON")
ap.add_argument("--week",   required=True, type=int)
ap.add_argument("--period", required=True, help='e.g. "04.05–10.05.2026"')
ap.add_argument("--today",  required=True, help='YYYY-MM-DD')
ap.add_argument("--extra-pages", type=int, default=0,
                help="Сколько дополнительных страниц будет приклеено к итоговому PDF (например, портретный слайд документов подрядчиков). Учитывается в нумерации N/M.")
ap.add_argument("--out-pptx", default=None)
ap.add_argument("--out-pdf-name", default=None,
                help='Override default PDF name; week is appended automatically')
args = ap.parse_args()

with open(args.config, encoding="utf-8") as f:
    CFG = json.load(f)

TPL           = CFG["template_pptx"]
COVER_IMG     = CFG["cover_image"]
OUT_DIR       = CFG["output_dir"]
PROJECT_NAME  = CFG["project_full"]
PROJECT_ADDR  = CFG["project_addr"]
PROJECT_SHORT = CFG["project_short"]
PROJECT_KEY   = CFG.get("name", "проект")
WEEK   = args.week
PERIOD = args.period
TODAY  = dt.datetime.strptime(args.today, "%Y-%m-%d").date()

REPORT_BASENAME = CFG.get("report_basename") or PROJECT_SHORT.replace('«','').replace('»','')
OUT_PPTX = args.out_pptx or os.path.join(OUT_DIR, f"Отчет из МСГ {REPORT_BASENAME}.pptx")
pdf_base = args.out_pdf_name or f"МСГ критические отставания {REPORT_BASENAME} неделя {WEEK}"
OUT_PDF  = os.path.join(OUT_DIR, f"{pdf_base}.pdf")

with open(args.items, encoding="utf-8") as f:
    ITEMS = json.load(f)

# ---------------------------- pct auto-normalize ----------------------------
# Авто-нормализация pct: некоторые проекты хранят процент готовности как ДОЛЮ
# (0..1), другие — как ПРОЦЕНТ (0..100). На Бугры-3 на 2026-05-21 pct хранится
# как доля (max=1.6, среднее 0.22), на Репино/Марьино — как проценты
# (Марьино max=1425, среднее 72.5). Без нормализации жизненный цикл показывает
# «0.6%/0.4%/...» вместо «60%/40%/...».
#
# Признак «pct в долях»: <5% значений превышают 1.5. Тогда умножаем все pct
# на 100. _to_pct ниже потом обрежет в [0,100], так что выбросы из обоих
# форматов (например 1425% у Марьино) корректно зажмутся.
def _parse_pct(v):
    if v is None: return None
    try:
        return float(str(v).replace(',', '.').replace('+', '').replace('%', '').replace(' ', '').replace('\xa0', '').strip())
    except Exception:
        return None

_pct_vals = [_parse_pct(it.get('pct')) for it in ITEMS]
_pct_vals = [v for v in _pct_vals if v is not None]
if _pct_vals:
    _above = sum(1 for v in _pct_vals if v > 1.5)
    if _above / len(_pct_vals) < 0.05:
        for _it in ITEMS:
            _f = _parse_pct(_it.get('pct'))
            if _f is not None:
                _it['pct'] = _f * 100.0
        print(f"INFO: pct нормализован x100 (был в долях: max={max(_pct_vals):.3f}, "
              f"вне-[0,1.5] {_above}/{len(_pct_vals)})")
    else:
        print(f"INFO: pct уже в процентах (max={max(_pct_vals):.1f}, "
              f"> 1.5: {_above}/{len(_pct_vals)})")

# ---------------------------- helpers ----------------------------
def _to_pct(v):
    if v is None: return None
    try:
        s = str(v).replace(',', '.').replace('+', '').replace('%', '').replace(' ', '').replace('\xa0', '').strip()
        f = float(s)
        if f > 100: f = 100.0
        if f < 0: f = 0.0
        return f
    except: return None

def _to_float(v):
    if isinstance(v, str):
        s = v.replace('\xa0', '').replace(' ', '').replace(',', '.').replace('+', '').strip()
        return float(s)
    return float(v)

def _to_float_safe(v):
    if v is None or v == '': return None
    try: return _to_float(v)
    except: return None

def parse_date(s):
    if not s: return None
    s = str(s).strip()
    if not s or s.startswith('#') or s == '-': return None
    for fmt in ("%Y-%m-%d", "%d.%m.%Y", "%d.%m.%y"):
        try:
            d = dt.datetime.strptime(s, fmt).date()
            if d.year < 2010: return None
            return d
        except: pass
    return None

def fmt_date(s):
    if not s: return '–'
    if isinstance(s, str) and (s.startswith('#') or s == '-'): return '–'
    d = parse_date(s)
    return d.strftime("%d.%m.%y") if d else '–'

def fmt_int(v):
    if v is None or v == '': return '–'
    try:
        n = int(v)
        return '0' if n == 0 else f"{n:+d}"
    except: return str(v)

def fmt_days_neg_only(v):
    """Показывать только отрицательные дни. При нуле или опережении — прочерк."""
    if v is None or v == '': return '–'
    try:
        n = int(v)
        if n >= 0: return '–'
        return f"{n:+d}"
    except: return str(v)

def fmt_num(v):
    if v is None or v == '': return '–'
    try:
        f = _to_float(v)
        if abs(f) < 0.005: return '0'
        if abs(f - round(f)) < 1e-6:
            return f"{int(round(f)):+d}"
        sign = '+' if f > 0 else ''
        return f"{sign}{f:.1f}".replace('.', ',')
    except: return str(v)

def fmt_pos(v):
    if v is None or v == '': return '–'
    try:
        f = _to_float(v)
        if abs(f - round(f)) < 1e-6:
            return str(int(round(f)))
        return f"{f:.1f}".replace('.', ',')
    except: return str(v)

# ---------------------------- lifecycle ----------------------------
def _avg_pct(rtypes):
    vals = [_to_pct(it['pct']) for it in ITEMS if it['rtype'] in rtypes]
    vals = [v for v in vals if v is not None]
    return sum(vals)/len(vals) if vals else 0.0

_stages = [
    (1, "Выпуск рабочей документации (РД)",          ['План РД'], ['Факт РД']),
    (2, "Пакет",                                      ['План П'],  ['Факт П']),
    (3, "Тендер",                                     ['План Т'],  ['Факт Т']),
    (4, "Заключение договора",                        ['План Д'],  ['Факт Д']),
    (5, "Мобилизация подрядчика и МТР",               ['План М'],  ['Факт М']),
    (6, "Строительно-монтажные работы (СМР)",         ['План'],    ['Факт']),
]
# NOTE: блок «Финансирование (оплата аванса)» временно отключён (2026-05-12).
# Полная версия с финансированием сохранена в build_report_v_with_finance.py
# на случай возврата позже.
LIFECYCLE = []
for n, name, pt, ft in _stages:
    p = _avg_pct(pt); f = _avg_pct(ft)
    LIFECYCLE.append((n, name, p, f, f - p))

print("Computed lifecycle:")
for n, nm, p, f, d in LIFECYCLE:
    print(f"  {n}. {nm}: План={p:.1f} Факт={f:.1f} Δ={d:+.1f}")

# ---------------------------- pair plan/fact rows ----------------------------
def build_pairs():
    pairs = []
    i = 0
    while i < len(ITEMS):
        cur = ITEMS[i]
        if i+1 < len(ITEMS) and cur['rtype'].startswith('План') and ITEMS[i+1]['rtype'].startswith('Факт'):
            pairs.append({'plan': cur, 'fact': ITEMS[i+1]})
            i += 2
        else:
            i += 1
    return pairs
PAIRS = build_pairs()

def is_completed(fact):
    return parse_date(fact['pf_nach']) is not None and parse_date(fact['pf_okon']) is not None

def lagging_by_stage(plan_rtype):
    out = []
    for p in PAIRS:
        if p['plan']['rtype'] != plan_rtype: continue
        plan, fact = p['plan'], p['fact']
        if is_completed(fact): continue
        n = _to_float_safe(plan.get('obj_otkl'))
        if n is None or n >= 0: continue
        plan_end = parse_date(plan['pf_okon'])
        fact_end = parse_date(fact['pf_okon'])
        days_dev = (plan_end - (fact_end or TODAY)).days if plan_end else None
        out.append({
            'name': plan['name'], 'block': plan.get('block') or '',
            'section': plan.get('section') or '',
            'contractor': plan.get('contractor') or '',
            'plan_n': plan['pf_nach'], 'plan_o': plan['pf_okon'],
            'fact_n': fact['pf_nach'], 'fact_o': fact['pf_okon'],
            'days_dev': days_dev,
            'obj_otkl': plan.get('obj_otkl'),
            'unit': plan.get('unit'),
            'pct': fact.get('pct'),
            'prereq': plan.get('prereq'),
        })
    return out

def lagging_smr():
    out = []
    for p in PAIRS:
        if p['plan']['rtype'] != 'План': continue
        plan, fact = p['plan'], p['fact']
        if is_completed(fact): continue
        n = _to_float_safe(plan.get('obj_otkl'))
        if n is None or n >= 0: continue
        plan_end = parse_date(plan['pf_okon'])
        fact_end = parse_date(fact['pf_okon'])
        days_dev = (plan_end - (fact_end or TODAY)).days if plan_end else None
        out.append({
            'name': plan['name'], 'block': plan.get('block') or '',
            'section': plan.get('section') or '',
            'contractor': plan.get('contractor') or '',
            'plan_n': plan['pf_nach'], 'plan_o': plan['pf_okon'],
            'fact_n': fact['pf_nach'], 'fact_o': fact['pf_okon'],
            'days_dev': days_dev,
            'obj_otkl': plan.get('obj_otkl'),
            'unit': plan.get('unit'),
            'proj_obj': plan.get('proj_obj'),
            'prereq': plan.get('prereq'),
        })
    return out

RD_LAG  = lagging_by_stage('План РД')
PAK_LAG = lagging_by_stage('План П')
TND_LAG = lagging_by_stage('План Т')
DOG_LAG = lagging_by_stage('План Д')
MOB_LAG = lagging_by_stage('План М')
SMR_LAG = lagging_smr()

print(f"Lag counts: РД={len(RD_LAG)} Пакет={len(PAK_LAG)} Тендер={len(TND_LAG)} "
      f"Договор={len(DOG_LAG)} Моб={len(MOB_LAG)} СМР={len(SMR_LAG)}")

# ---------------------------- shape helpers ----------------------------
prs = Presentation(TPL)

# Высота строк таблиц критических отставаний:
# HEADER_ROW_H — строка-«оглавление» (заголовки колонок), оставляем повыше.
# UNIFORM_ROW_H — секционные разделители + строки данных, чуть компактнее.
HEADER_ROW_H = 240000   # EMU, ~25pt
UNIFORM_ROW_H = 180000  # EMU, ~18.7pt

def set_text(shape, new_text):
    if not shape.has_text_frame: return
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        first_run = p.runs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        new_r = p.add_run(); new_r.text = new_text
        try:
            new_r.font.name = first_run.font.name
            new_r.font.size = first_run.font.size
            new_r.font.bold = first_run.font.bold
            if first_run.font.color and first_run.font.color.type is not None:
                new_r.font.color.rgb = first_run.font.color.rgb
        except: pass
    else:
        new_r = p.add_run(); new_r.text = new_text
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)

def remove_shape(shape):
    shape._element.getparent().remove(shape._element)

def clone_slide(prs, src_slide):
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)
    spTree = new_slide.shapes._spTree
    for shape in src_slide.shapes:
        spTree.append(deepcopy(shape._element))
    return new_slide

SLIDES_TO_REMOVE = []  # collect slides to drop AFTER all clones are done

def remove_slide(prs, slide):
    """Remove a slide from the deck. Important: when called BEFORE any
    subsequent clone_slide() calls, python-pptx's next_partname allocator
    can reuse the freed partname and produce duplicate ZIP entries on save.
    To avoid that, defer removal: append to SLIDES_TO_REMOVE and call
    _flush_slide_removals(prs) once near the end."""
    SLIDES_TO_REMOVE.append(slide)

def _flush_slide_removals(prs):
    sldIdLst = prs.slides._sldIdLst
    for slide in SLIDES_TO_REMOVE:
        target_rId = None
        for rId, rel in prs.part.rels.items():
            if rel.target_part is slide.part:
                target_rId = rId; break
        for el in list(sldIdLst):
            if el.get(qn('r:id')) == target_rId:
                sldIdLst.remove(el); break
    SLIDES_TO_REMOVE.clear()

def reorder_slide_to(prs, slide, target_idx):
    sldIdLst = prs.slides._sldIdLst
    target_rId = None
    for rId, rel in prs.part.rels.items():
        if rel.target_part is slide.part:
            target_rId = rId; break
    for el in list(sldIdLst):
        if el.get(qn('r:id')) == target_rId:
            sldIdLst.remove(el)
            sldIdLst.insert(target_idx, el)
            return

# ---------------------------- slide 1: cover ----------------------------
slide1 = prs.slides[0]
# project name + address — большой заголовок наверху, явно НАД банером МСГ
for shape in slide1.shapes:
    if shape.has_text_frame and 'Репин' in shape.text_frame.text:
        tf = shape.text_frame
        shape.left = Emu(0); shape.top = Emu(300000)
        shape.width = prs.slide_width; shape.height = Emu(1700000)
        tf.word_wrap = True
        try: tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception: pass
        for p in list(tf.paragraphs):
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = PROJECT_NAME
        r1.font.bold = True; r1.font.size = Pt(40); r1.font.name = 'Calibri'
        r1.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        p2 = tf.paragraphs[1] if len(tf.paragraphs) >= 2 else tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = PROJECT_ADDR
        r2.font.size = Pt(18); r2.font.name = 'Calibri'
        r2.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        for extra in tf.paragraphs[2:]:
            extra._p.getparent().remove(extra._p)
        break

# big banner (title + week) — 2 lines: «ОТЧЁТ» / «НЕДЕЛЯ XX (...)»
# Совещание с ОЛ от 19.05.2026: убрать аббревиатуру «МСГ» и подзаголовок
# «КРИТИЧЕСКИЕ ОТСТАВАНИЯ» — в отчёте уже не только МСГ-данные, и не только
# отставания (есть и недоработки по документам).
for shape in slide1.shapes:
    if shape.has_text_frame and 'НЕДЕЛЯ' in shape.text_frame.text:
        tf = shape.text_frame
        shape.left = Emu(0); shape.top = Emu(2700000)
        shape.width = prs.slide_width; shape.height = Emu(2400000)
        for p in list(tf.paragraphs):
            p._p.getparent().remove(p._p)
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = "ОТЧЁТ"
        r1.font.bold = True; r1.font.size = Pt(72); r1.font.name = 'Calibri'
        r1.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = f"НЕДЕЛЯ {WEEK} ({PERIOD})"
        r2.font.bold = True; r2.font.size = Pt(40); r2.font.name = 'Calibri'
        r2.font.color.rgb = RGBColor(0xE6, 0x00, 0x73)
        break

# cover image: full-slide background, 20% opacity, behind text
for shape in list(slide1.shapes):
    if shape.shape_type == 13 and shape.name == 'Рисунок 4':
        remove_shape(shape); break
bg_pic = slide1.shapes.add_picture(COVER_IMG, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
blip_fill = bg_pic._element.find('.//' + qn('p:blipFill'))
if blip_fill is not None:
    blip = blip_fill.find(qn('a:blip'))
    if blip is not None:
        alpha = etree.SubElement(blip, qn('a:alphaModFix'))
        alpha.set('amt', '20000')
spTree = bg_pic._element.getparent()
spTree.remove(bg_pic._element)
insert_idx = 0
for i, child in enumerate(spTree):
    tag = etree.QName(child).localname
    if tag in ('nvGrpSpPr', 'grpSpPr'):
        insert_idx = i + 1
spTree.insert(insert_idx, bg_pic._element)

# Pre-capture ALL slide refs BEFORE any insertions / clones / removals so
# subsequent index shifts (e.g. inserting the ГПР slide at position 1) don't
# break later references. Use object refs below, not indices.
slide2 = prs.slides[1]  # lifecycle
slide3 = prs.slides[2]  # СМР
slide4 = prs.slides[3]  # РД
slide5 = prs.slides[4]  # Тендер
slide6 = prs.slides[5]  # Договор
slide7 = prs.slides[6]  # Финансирование
slide8 = prs.slides[7]  # Мобилизация

# ---------------------------- slide 2 (insert at idx 1): ГПР readiness screenshot ----------------------------
# Done EARLY (before any clone_slide / remove_slide for the stage tables) so
# python-pptx's partname allocator gives the new slide a fresh number and
# doesn't collide with cloned/removed slide partnames.
# Source: single sheet «Отчеты Март 2026» with 4 object blocks side-by-side.
# Project is identified by config["gpr_object_match"] matched against the
# 'Отчет о готовности ГПР и МСГ по объекту "<name>" ...' header cell.
GPR_REPORTS_DIR = r"C:\Авраменко\1. КОМПАКТ\Отчеты"
GPR_OBJECT_MATCH = CFG.get("gpr_object_match")

def _find_latest_gpr_xlsx():
    import glob, re
    pattern = os.path.join(GPR_REPORTS_DIR, "Отчет о готовности ГПР и МСГ по объектам*.xlsx")
    candidates = [p for p in glob.glob(pattern) if not os.path.basename(p).startswith("~$")]
    if not candidates: return None
    def keyof(p):
        m = re.search(r'от (\d{2})\.(\d{2})\.(\d{4})', os.path.basename(p))
        if m: return (int(m.group(3)), int(m.group(2)), int(m.group(1)))
        return (0, 0, 0)
    return max(candidates, key=keyof)

def _gpr_date_from_name(p):
    import re
    m = re.search(r'от (\d{2}\.\d{2}\.\d{4})', os.path.basename(p))
    return m.group(1) if m else ""

def _gpr_locate_block(xlsx_path, keyword):
    """Return (sheet_name, start_row, end_row) for the requested object."""
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, data_only=True, read_only=False)
    try:
        for sn in wb.sheetnames:
            ws = wb[sn]
            for row in ws.iter_rows(min_row=1, max_row=ws.max_row, values_only=False):
                for cell in row:
                    v = cell.value
                    if isinstance(v, str) and "Отчет о готовности" in v and keyword in v:
                        start = cell.row
                        # block ends at the first fully-empty row in cols C..J after the header
                        end = start + 24
                        for r in range(start + 3, min(start + 35, ws.max_row + 1)):
                            empty = all(
                                (ws.cell(r, c).value is None or str(ws.cell(r, c).value).strip() == "")
                                for c in range(3, 11)
                            )
                            if empty:
                                end = r - 1
                                break
                        return sn, start, end
        return None
    finally:
        wb.close()

def _gpr_has_issues(xlsx_path, sheet_name, start_row, end_row):
    """ОЛ 19.05.2026: ГПР-слайд показываем только если есть проблемы.
    Считаем проблемой: F='не готов' или G<1.0 (есть, но не дотягивает до 100%).
    Пустые значения — это «не планируется» или «в работе», не считаем
    проблемой. Если в блоке нет ни одной проблемной строки — слайд пропускаем.
    """
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, data_only=True, read_only=False)
    try:
        ws = wb[sheet_name]
        for r in range(start_row + 1, end_row + 1):
            f_val = ws.cell(r, 6).value  # «Отметка о заполнении»
            g_val = ws.cell(r, 7).value  # «% готовности»
            if isinstance(f_val, str) and "не готов" in f_val.lower():
                return True
            if isinstance(g_val, (int, float)) and 0 < g_val < 1.0:
                return True
        return False
    finally:
        wb.close()


def _gpr_screenshot(xlsx_path, sheet_name, range_addr, out_png):
    import time, pythoncom, win32com.client
    from PIL import ImageGrab
    pythoncom.CoInitialize()
    excel = win32com.client.DispatchEx("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False
    try:
        wb = excel.Workbooks.Open(xlsx_path, ReadOnly=True)
        ws = wb.Sheets(sheet_name)
        ws.Activate()
        # Force Normal view — Page Break Preview overlays a "Страница N" watermark
        try: excel.ActiveWindow.View = 1  # xlNormalView
        except Exception: pass
        rng = ws.Range(range_addr)
        rng.CopyPicture(Appearance=1, Format=2)  # xlScreen=1, xlBitmap=2
        time.sleep(0.6)
        img = ImageGrab.grabclipboard()
        if img is None:
            raise RuntimeError("Excel CopyPicture produced no clipboard image")
        img.save(out_png, "PNG")
        wb.Close(SaveChanges=False)
    finally:
        excel.Quit()
        pythoncom.CoUninitialize()

# ---------------------------- Предписания: download + build summary table ----------------------------
def _bitrix_download_disk_file(file_id, out_path):
    """Скачать файл с Битрикс-Диска по его file_id через REST disk.file.get."""
    import urllib.request, json as _json
    webhook_path = r"C:\Авраменко\Claude Code Projects\.bitrix-webhook"
    webhook = ""
    with open(webhook_path, encoding="utf-8") as f:
        for line in f:
            s = line.strip()
            if not s or s.startswith("#"):
                continue
            webhook = s.split()[0]
            break
    if not webhook:
        raise RuntimeError(f"webhook is empty: {webhook_path}")
    if not webhook.startswith("http"):
        webhook = f"https://kspb.bitrix24.ru/rest/2958/{webhook}/"
    if not webhook.endswith("/"):
        webhook += "/"
    info_url = f"{webhook}disk.file.get?id={file_id}"
    # Битрикс/Google периодически рвут коннект — делаем до 4 попыток на каждый
    # шаг (получение DOWNLOAD_URL + сам download).
    import time as _t
    _last = None
    for _attempt in range(1, 5):
        try:
            with urllib.request.urlopen(info_url, timeout=30) as r:
                data = _json.loads(r.read().decode("utf-8"))
            dl = data["result"]["DOWNLOAD_URL"]
            urllib.request.urlretrieve(dl, out_path)
            return out_path
        except Exception as _e:
            _last = _e
            print(f"  bitrix download attempt {_attempt}/4 failed: {_e}")
            _t.sleep(2)
    raise _last

def _pred_load_filtered_rows(src_xlsx):
    """Читает сводную таблицу предписаний СК и возвращает список 8-кортежей
    по строкам, у которых K='-' (СК ещё не снял предписание).

    Структура файла (новая, с 19.05.2026 — «Сводная таблица предписаний СК
    (Репино).xlsx» file_id 1639646; шапка в R3, данные с R4):
      B №предп · C дата выдачи · D компания · E перечень нарушений ·
      F срок устранения · G Срыв срока (число) · H Должность Ф.И.О. ·
      I Статус выполнения · J Отметка об устранении/прогноз ·
      K Отметка о снятии СК ('+'/'-' или пусто).

    Merge-родители B/C/D подставляются из last-seen. Title-строки (только E
    заполнен) пропускаются. Фильтр K=='-' — показываем только не снятые СК
    позиции."""
    import openpyxl as _ox
    import datetime as _dt
    wb_s = _ox.load_workbook(src_xlsx, data_only=True)
    ws_s = wb_s.active
    def _fmt(v):
        if isinstance(v, (_dt.datetime, _dt.date)):
            return v.strftime("%d.%m.%Y")
        if v is None: return ""
        if isinstance(v, float) and v.is_integer():
            return str(int(v))
        return str(v).strip()
    last_B = last_C = last_D = None
    # Заголовки разделов («Предписания строительного контроля АО СК "Компакт"»,
    # «Предписания строительного контроля Заказчика») сейчас занимают title-
    # строки, у которых заполнен только E. Раньше мы их пропускали — теперь
    # они идут в выход как «section»-маркеры ПЕРЕД первой data-строкой своего
    # раздела (только если под ними есть хоть одна не-снятая позиция).
    pending_section = None
    rows = []
    for r in range(4, ws_s.max_row + 1):
        e = ws_s.cell(row=r, column=5).value
        cells = [ws_s.cell(row=r, column=c).value for c in range(2, 12)]
        if e and all(cells[i] is None for i in (0,1,2,4,5,6,7,8,9)):
            # Title-строка раздела — запоминаем, но НЕ emit'им пока не появится
            # data-строка под ней. Сбрасывает merge-родителей B/C/D — при смене
            # раздела last_B/C/D не должны протекать.
            pending_section = str(e).strip()
            last_B = last_C = last_D = None
            continue
        b = ws_s.cell(row=r, column=2).value
        c2 = ws_s.cell(row=r, column=3).value
        d = ws_s.cell(row=r, column=4).value
        if b is not None: last_B = b
        if c2 is not None: last_C = c2
        if d is not None: last_D = d
        # K — «Отметка о снятии СК». Берём только не снятые ('-').
        k_snyatie = ws_s.cell(row=r, column=11).value
        if (str(k_snyatie).strip() if k_snyatie is not None else "") != "-":
            continue
        if pending_section is not None:
            rows.append({"section": pending_section})
            pending_section = None
        j_val = _fmt(ws_s.cell(row=r, column=10).value)
        rows.append((
            _fmt(last_B), _fmt(last_C), _fmt(last_D),
            _fmt(ws_s.cell(row=r, column=5).value),  # E: Нарушения
            _fmt(ws_s.cell(row=r, column=6).value),  # F: Срок устранения
            _fmt(ws_s.cell(row=r, column=7).value),  # G: Срыв срока (число)
            j_val or "–",                            # J: Отметка/прогноз
            "Не снято",                              # K='-' → не снято СК
        ))
    wb_s.close()
    return rows

if GPR_OBJECT_MATCH:
    gpr_xlsx = _find_latest_gpr_xlsx()
    if gpr_xlsx:
        loc = _gpr_locate_block(gpr_xlsx, GPR_OBJECT_MATCH)
        _gpr_skipped_ok = False
        if loc and not _gpr_has_issues(gpr_xlsx, loc[0], loc[1], loc[2]):
            print(f"INFO: ГПР по '{GPR_OBJECT_MATCH}' — проблем нет, слайд пропущен "
                  f"(правка ОЛ 19.05.2026)")
            loc = None
            _gpr_skipped_ok = True
        if loc:
            gpr_sheet_name, gpr_r1, gpr_r2 = loc
            # Skip the top "Отчет о готовности…" caption row (gpr_r1). Keep the
            # blank spacer row (gpr_r1+1) as a top margin so the upper border of
            # the column-headers row isn't clipped at the screenshot edge.
            # Include empty column A as a left margin for the same reason.
            gpr_range = f"A{gpr_r1+1}:I{gpr_r2}"
            import tempfile
            tmp_png = os.path.join(tempfile.gettempdir(), f"gpr_{PROJECT_KEY}.png")
            try:
                _gpr_screenshot(gpr_xlsx, gpr_sheet_name, gpr_range, tmp_png)
                # Clone slide3 (СМР) — at this point it's still the pristine
                # template-state and has a header bar identical to the rest.
                gpr_slide = clone_slide(prs, slide3)
                # Strip non-header shapes: keep separator line, title rectangle,
                # slide-number rectangle, КОМПАКТ logo. Drop everything else.
                _HEADER_NAMES = {
                    'Прямая соединительная линия 10',
                    'Прямоугольник 11',
                    'Прямоугольник 12',
                    'Рисунок 6',
                }
                for _sh in list(gpr_slide.shapes):
                    if getattr(_sh, 'name', '') not in _HEADER_NAMES:
                        remove_shape(_sh)
                # Update title text. ОЛ 19.05.2026: «готовность» → «актуализация» —
                # лист по факту фиксирует, актуализируются данные или нет, а не
                # измеряет «готовность» как таковую.
                gpr_title = f"АКТУАЛИЗАЦИЯ ГПР И МСГ ({PROJECT_SHORT})"
                for _sh in gpr_slide.shapes:
                    if getattr(_sh, 'name', '') == 'Прямоугольник 11':
                        set_text(_sh, gpr_title); break
                # Add screenshot under the header
                sw, sh = prs.slide_width, prs.slide_height
                from PIL import Image
                with Image.open(tmp_png) as _im:
                    _iw, _ih = _im.size
                _ml = Emu(341998)
                _aw = sw - _ml * 2
                _top = Emu(900000)
                _ah = sh - _top - Emu(300000)
                _scale = min(_aw / _iw, _ah / _ih)
                _nw = int(_iw * _scale); _nh = int(_ih * _scale)
                _pl = (sw - _nw) // 2
                _pt = _top + (_ah - _nh) // 2
                gpr_slide.shapes.add_picture(tmp_png, _pl, _pt, width=_nw, height=_nh)
                # Move to position 1 (right after cover)
                reorder_slide_to(prs, gpr_slide, 1)
                print(f"GPR slide inserted: {os.path.basename(gpr_xlsx)} / "
                      f"{gpr_sheet_name}!{gpr_range} (match='{GPR_OBJECT_MATCH}')")
            except Exception as _e:
                print(f"WARN: GPR slide skipped: {_e}")
        elif not _gpr_skipped_ok:
            print(f"WARN: object '{GPR_OBJECT_MATCH}' not found in {os.path.basename(gpr_xlsx)}")
    else:
        print(f"WARN: no ГПР xlsx found in {GPR_REPORTS_DIR}")
else:
    print("INFO: gpr_object_match not set in config — skipping ГПР slide")

# ---------------------------- slide 2: lifecycle table ----------------------------
for s in slide2.shapes:
    if s.has_text_frame and 'СВОДНАЯ' in s.text_frame.text:
        set_text(s, f"ЖИЗНЕННЫЙ ЦИКЛ МСГ ({PROJECT_SHORT})")
        break
for s in list(slide2.shapes):
    if s.shape_type == 13 and s.name == 'Рисунок 4':
        # Правка пользователя 20.05.2026: ширина таблицы — на всю ширину слайда
        # как у стадийных таблиц (W = slide_w − 2·341998 EMU). Y/высоту
        # наследуем от шейпа-плейсхолдера «Рисунок 4».
        _, T, _, H = s.left, s.top, s.width, s.height
        L = Emu(341998)
        W = prs.slide_width - Emu(341998) * 2
        remove_shape(s)
        table_shape = slide2.shapes.add_table(rows=len(LIFECYCLE)+1, cols=5, left=L, top=T, width=W, height=H)
        tbl = table_shape.table
        for ci, prop in enumerate([0.06, 0.50, 0.14, 0.14, 0.16]):
            tbl.columns[ci].width = int(W * prop)
        # Сводная — строки чуть выше стадийных (всего 7 строк, есть запас)
        LIFECYCLE_ROW_H = 360000  # EMU ≈ 37.5pt
        for _r in tbl.rows:
            _r.height = LIFECYCLE_ROW_H
        for _r in tbl.rows:
            for _c in _r.cells:
                _c.margin_top = Emu(25000); _c.margin_bottom = Emu(25000)
                _c.margin_left = Emu(60000); _c.margin_right = Emu(60000)
                _c.text_frame.word_wrap = True
        for ci, h in enumerate(["№", "Этап жизненного цикла", "План, %", "Факт, %", "Δ (Факт−План), п.п."]):
            cell = tbl.cell(0, ci); cell.text = ""
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = h
            r.font.bold = True; r.font.size = Pt(14); r.font.name = 'Calibri'
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for ri, row in enumerate(LIFECYCLE, start=1):
            num, name, plan, fact, delta = row
            # ОЛ 19.05.2026: при Факт = 100 % подсвечиваем строку зелёным —
            # «строка открывается / стадия закрыта».
            is_complete = fact >= 100
            vals = [str(num), name, f"{plan:.1f}".replace('.', ','), f"{fact:.1f}".replace('.', ','),
                    f"{'+' if delta>0 else ''}{delta:.1f}".replace('.', ',')]
            for ci, v in enumerate(vals):
                cell = tbl.cell(ri, ci); cell.text = ""
                if is_complete:
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xC6, 0xEF, 0xCE)
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if ci == 1 else PP_ALIGN.CENTER
                r = p.add_run(); r.text = v
                r.font.size = Pt(13); r.font.name = 'Calibri'
                if is_complete:
                    r.font.bold = True
                if ci == 4:
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(0xE0,0,0) if delta < 0 else RGBColor(0,0x80,0)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        break

# ---------------------------- slide 2b: Численность подрядчиков (только Репино) ----------------------------
# Источник — лист `Численность` в gsheet `МСГ_RBI Репино Санаторий` (тот же
# gdrive_file_id, что для МСГ).
# ОЛ 19.05.2026: ушли от показа 7 дней недели. Теперь по каждому подрядчику ×
# категория показываем: «Пред. неделя» / «Тек. неделя» / Δ (тек − пред). Это
# сравнение динамики двух точек, без ежедневной развёртки. Норматив (НТД)
# планируется отдельной задачей — сейчас столбец «критическое отклонение»
# не выводится. Подключается флагом include_chislennost в config; вставляется
# СРАЗУ ПОСЛЕ слайда «Жизненный цикл», ПЕРЕД divider'ом СМР.
CHISL_SLIDE = None
if CFG.get("include_chislennost"):
    try:
        import tempfile as _tf
        import urllib.request as _urlreq
        import openpyxl as _ox

        _RU_MONTHS_SHORT = {1:"Янв",2:"Фев",3:"Мар",4:"Апр",5:"Май",6:"Июн",
                            7:"Июл",8:"Авг",9:"Сен",10:"Окт",11:"Ноя",12:"Дек"}

        _chisl_xlsx = os.path.join(_tf.gettempdir(), f"chisl_{PROJECT_KEY}.xlsx")
        _gsheet_url = f"https://docs.google.com/spreadsheets/d/{CFG['gdrive_file_id']}/export?format=xlsx"
        print(f"Численность: GET {_gsheet_url}")
        # Google Docs нестабилен — рвёт коннект на больших экспортах. Делаем
        # до 4 попыток с короткими паузами.
        import time as _time
        _last_exc = None
        for _attempt in range(1, 5):
            try:
                _urlreq.urlretrieve(_gsheet_url, _chisl_xlsx)
                _last_exc = None
                break
            except Exception as _e_retry:
                _last_exc = _e_retry
                print(f"  попытка {_attempt}/4 не удалась: {_e_retry}")
                _time.sleep(2)
        if _last_exc is not None:
            raise _last_exc
        _wb_c = _ox.load_workbook(_chisl_xlsx, data_only=True)
        # Лист переименовывали: «Численность» → «Люди, техника» → «Ресурсы»
        # (20.05.2026, унифицировано по обоим проектам — Репино и Марьино).
        # Структура та же — поддерживаем все три варианта.
        _CHISL_SHEET_CANDIDATES = ("Ресурсы", "Численность", "Люди, техника")
        _chisl_sheet_name = next((_n for _n in _CHISL_SHEET_CANDIDATES
                                  if _n in _wb_c.sheetnames), None)
        if _chisl_sheet_name is None:
            raise RuntimeError(f"лист численности не найден в gsheet "
                               f"(пробовали: {_CHISL_SHEET_CANDIDATES}, "
                               f"в файле: {_wb_c.sheetnames})")
        _ws_c = _wb_c[_chisl_sheet_name]

        # Кеш {(год, месяц): min_col} — первый день месяца в шапке R1 (merge-блок).
        _month_first_col = {}
        for _mc in _ws_c.merged_cells.ranges:
            if _mc.min_row == 1 and _mc.max_row == 1:
                _tl = _ws_c.cell(row=1, column=_mc.min_col).value
                if not isinstance(_tl, str): continue
                _parts = _tl.strip().split()
                if len(_parts) != 2: continue
                _mname, _myear = _parts
                for _mnum, _mlab in _RU_MONTHS_SHORT.items():
                    if _mlab == _mname:
                        try:
                            _month_first_col[(int(_myear), _mnum)] = _mc.min_col
                        except ValueError:
                            pass
                        break

        def _col_for_date(d):
            """Колонка в листе «Численность», соответствующая дате d. None если месяц не найден."""
            base = _month_first_col.get((d.year, d.month))
            if base is None: return None
            return base + (d.day - 1)

        # Отчётный день — СРЕДА (правка 20.05.2026). Раньше брали понедельник,
        # но руководство сравнивает динамику «по средам»: тек.ср = ср недели
        # пятничной публикации (= TODAY если TODAY — среда, иначе ближайшая
        # прошедшая), пред.ср = ср предыдущей недели. weekday()==2 → среда.
        _curr_day = TODAY - dt.timedelta(days=(TODAY.weekday() - 2) % 7)
        _prev_day = _curr_day - dt.timedelta(days=7)
        _curr_col = _col_for_date(_curr_day)
        _prev_col = _col_for_date(_prev_day)
        print(f"Ресурсы: тек.ср {_curr_day.strftime('%d.%m.%Y')}, "
              f"пред.ср {_prev_day.strftime('%d.%m.%Y')}")

        # 20.05.2026: вместо хардкоженного списка подрядчиков сканируем колонку A
        # (имя\nроль). Это позволяет одной кодой работать и Репино и Марьино
        # (где подрядчики разные, ~14 шт.). Стартуем с R4, шаг 4 (4 категории).
        # Останавливаемся на пустом A или встретив «ИТОГО» (это тотал-блок).
        _CHISL_LAYOUT = []
        _r = 4
        while True:
            _a = _ws_c.cell(row=_r, column=1).value
            if not isinstance(_a, str) or not _a.strip():
                break
            if 'итого' in _a.lower().strip()[:8]:
                break
            _parts = _a.split('\n', 1)
            _name_p = _parts[0].strip()
            _role_p = _parts[1].strip() if len(_parts) > 1 else ''
            _CHISL_LAYOUT.append((_r, _name_p, _role_p))
            _r += 4
        print(f"Численность: автоопределено {len(_CHISL_LAYOUT)} подрядчиков "
              f"(R4..R{_r - 4})")
        # 20.05.2026: по запросу пользователя поменяли местами ИТР и Σ — теперь
        # Σ (общая численность) идёт первой строкой блока подрядчика, ИТР третьей.
        # Это меняет только визуальный порядок строк; данные пересобираются ниже
        # в _cats_vals/_total_cats так, чтобы соответствовать новому _CATS.
        _CATS = ["Σ", "Раб", "ИТР", "Техн."]  # 4 категории на подрядчика (offsets 0..3)

        def _chisl_get(row, col):
            if col is None: return None
            v = _ws_c.cell(row=row, column=col).value
            if v is None or v == "": return None
            try:
                return float(v)
            except (TypeError, ValueError):
                return None

        def _safe_sum(*vals):
            """Сумма, игнорируя None. Если все None — None."""
            _nums = [_v for _v in vals if _v is not None]
            return sum(_nums) if _nums else None

        # Для каждого подрядчика — 4 категории × {пред, тек}.
        # Σ-строка gsheet — это формула SUM от пустых = 0, поэтому Σ считаем
        # сами как ИТР + Раб (а не берём из gsheet).
        # Подсчёт значений + фильтр «пустых» подрядчиков (20.05.2026): если ВСЕ
        # 6 сырых значений (ИТР пред/тек, Раб пред/тек, Техн пред/тек) пустые —
        # подрядчик не активен на этих неделях, в слайд не идёт. Заменяет ручное
        # скрытие Болверк-Норда для Репино (работы по шпунту закончены).
        # Fallback: если ВСЕ подрядчики пустые (лист новый, planner ещё не
        # заполнил) — фильтр отключаем и показываем всех, иначе слайд был бы
        # пустым.
        _candidates = []
        for _r, _name, _role in _CHISL_LAYOUT:
            _v_itr_p  = _chisl_get(_r + 0, _prev_col)
            _v_itr_c  = _chisl_get(_r + 0, _curr_col)
            _v_rab_p  = _chisl_get(_r + 1, _prev_col)
            _v_rab_c  = _chisl_get(_r + 1, _curr_col)
            _v_thn_p  = _chisl_get(_r + 3, _prev_col)
            _v_thn_c  = _chisl_get(_r + 3, _curr_col)
            _v_sig_p  = _safe_sum(_v_itr_p, _v_rab_p)
            _v_sig_c  = _safe_sum(_v_itr_c, _v_rab_c)
            _is_empty = all(_v is None for _v in (_v_itr_p, _v_itr_c, _v_rab_p, _v_rab_c,
                                                  _v_thn_p, _v_thn_c))
            # Порядок строго соответствует _CATS: Σ / Раб / ИТР / Техн.
            _cats_vals = [
                (_v_sig_p, _v_sig_c),
                (_v_rab_p, _v_rab_c),
                (_v_itr_p, _v_itr_c),
                (_v_thn_p, _v_thn_c),
            ]
            _candidates.append((_name, _role, _cats_vals, _is_empty))
        _has_any_data = any(not c[3] for c in _candidates)
        _chisl_rows = []
        for _name, _role, _cats_vals, _is_empty in _candidates:
            if _has_any_data and _is_empty:
                print(f"  пропуск (нет данных): {_name}")
                continue
            _chisl_rows.append((len(_chisl_rows) + 1, _name, _role, _cats_vals))

        # ИТОГО: сумма по всем подрядчикам для каждой категории. Σ ИТОГО =
        # сумма Σ подрядчиков (== сумма ИТР + сумма Раб). Из gsheet не берём.
        _total_cats = []
        for _off in range(4):
            _v_prev = _safe_sum(*(row[3][_off][0] for row in _chisl_rows))
            _v_curr = _safe_sum(*(row[3][_off][1] for row in _chisl_rows))
            _total_cats.append((_v_prev, _v_curr))
        print(f"Численность: подрядчиков {len(_chisl_rows)}, "
              f"ИТОГО Σ пред={_total_cats[2][0]} тек={_total_cats[2][1]}")

        # Общие константы/хелперы для всех Ресурсы-слайдов (пагинация ниже).
        _HEADER_NAMES_CHISL = {
            'Прямая соединительная линия 10',
            'Прямоугольник 11',
            'Прямоугольник 12',
            'Рисунок 6',
        }
        _NAVY        = RGBColor(0x1F, 0x4E, 0x79)
        _WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
        _BLACK       = RGBColor(0x00, 0x00, 0x00)
        _ROLE        = RGBColor(0x60, 0x60, 0x60)
        _TOTAL_FILL  = RGBColor(0xDC, 0xE6, 0xF1)  # ИТОГО (4 нижние строки)
        _SIGMA_FILL  = RGBColor(0xEA, 0xF1, 0xF8)  # светло-голубой — теперь подсвечивает ИТР
        _RED_TXT     = RGBColor(0xC0, 0x00, 0x00)
        _GREEN_TXT   = RGBColor(0x00, 0x80, 0x00)
        # Ширины: №(0.04), Подрядчик(0.42), Кат.(0.08), Пред.(0.15), Тек.(0.15), Δ(0.16)
        _CHISL_FRACS = [0.04, 0.42, 0.08, 0.15, 0.15, 0.16]

        def _chisl_fmt_int(v):
            if v is None: return "–"
            try:
                return str(int(round(float(v))))
            except (TypeError, ValueError):
                return "–"

        def _chisl_set_simple(cell, text, *, bold=False, size=10, align=PP_ALIGN.CENTER,
                              fill=None, color=None, italic=False):
            if color is None:
                color = _WHITE if (fill is not None and fill == _NAVY) else _BLACK
            cell.text = ""
            p = cell.text_frame.paragraphs[0]; p.alignment = align
            r = p.add_run(); r.text = str(text)
            r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
            cell.text_frame.word_wrap = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Emu(5000); cell.margin_bottom = Emu(5000)
            cell.margin_left = Emu(30000); cell.margin_right = Emu(30000)
            if fill is not None:
                cell.fill.solid(); cell.fill.fore_color.rgb = fill

        def _chisl_delta_text(v_prev, v_curr):
            if v_prev is None or v_curr is None:
                return ("–", None)
            d = int(round(float(v_curr) - float(v_prev)))
            if d == 0:
                return ("0", None)
            if d > 0:
                return (f"+{d}", _GREEN_TXT)
            return (f"{d}", _RED_TXT)

        def _build_chisl_table(_slide, _rows_chunk, _include_total):
            """Собрать таблицу Ресурсов на одном слайде. _rows_chunk — список
            кортежей (idx, name, role, cats_vals) только для этого слайда;
            ИТОГО (4 строки) добавляется лишь при _include_total=True."""
            _n_rows = 1 + len(_rows_chunk) * 4 + (4 if _include_total else 0)
            _n_cols = 6
            _sw_e, _sh_e = prs.slide_width, prs.slide_height
            _W = _sw_e - Emu(341998) * 2
            _L = Emu(341998); _T = Emu(900000)
            _H = _sh_e - _T - Emu(900000)
            _tbl_shape = _slide.shapes.add_table(rows=_n_rows, cols=_n_cols,
                                                  left=_L, top=_T, width=_W, height=_H)
            _tbl = _tbl_shape.table
            for _ci, _fr in enumerate(_CHISL_FRACS):
                _tbl.columns[_ci].width = int(_W * _fr)
            _HDR_H = 360000
            # Высота data-строки фиксируется по «полной» странице (max подрядчиков
            # без ИТОГО). На неполном последнем чанке (Марьино: 5+ИТОГО) строки
            # НЕ растягиваются — высота как на предыдущей полной странице.
            # Правка пользователя 2026-05-20.
            _FULL_ROWS = 1 + _CHISL_PER_SLIDE * 4
            _DATA_H = int((_H - _HDR_H) / max(1, _FULL_ROWS - 1))
            _tbl.rows[0].height = Emu(_HDR_H)
            for _r in range(1, _n_rows):
                _tbl.rows[_r].height = Emu(_DATA_H)

            # Шапка
            _prev_lbl = f"Пред.\n{_prev_day.strftime('%d.%m')}"
            _curr_lbl = f"Тек.\n{_curr_day.strftime('%d.%m')}"
            _headers = ["№", "Подрядчик", "Кат.", _prev_lbl, _curr_lbl, "Δ\n(тек−пред)"]
            for _ci, _h in enumerate(_headers):
                _cell = _tbl.cell(0, _ci); _cell.text = ""
                _cell.fill.solid(); _cell.fill.fore_color.rgb = _NAVY
                _cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                _cell.margin_top = Emu(20000); _cell.margin_bottom = Emu(20000)
                _cell.margin_left = Emu(20000); _cell.margin_right = Emu(20000)
                _tf_h = _cell.text_frame; _tf_h.word_wrap = True
                for _p in list(_tf_h.paragraphs)[1:]:
                    _p._p.getparent().remove(_p._p)
                _lines = _h.split("\n")
                _p0 = _tf_h.paragraphs[0]; _p0.alignment = PP_ALIGN.CENTER
                _r0 = _p0.add_run(); _r0.text = _lines[0]
                _r0.font.name = "Calibri"; _r0.font.size = Pt(11)
                _r0.font.bold = True; _r0.font.color.rgb = _WHITE
                for _ln in _lines[1:]:
                    _p_ext = _tf_h.add_paragraph(); _p_ext.alignment = PP_ALIGN.CENTER
                    _r_ext = _p_ext.add_run(); _r_ext.text = _ln
                    _r_ext.font.name = "Calibri"; _r_ext.font.size = Pt(9)
                    _r_ext.font.bold = False; _r_ext.font.color.rgb = _WHITE

            # Строки данных
            for _pi, (_idx, _name, _role, _cats_vals) in enumerate(_rows_chunk):
                _row_start = 1 + _pi * 4
                # № — merge 4 строк
                _cell_num = _tbl.cell(_row_start, 0)
                _cell_num.merge(_tbl.cell(_row_start + 3, 0))
                _chisl_set_simple(_cell_num, str(_idx), size=12, bold=True)
                # Подрядчик — merge 4 строк
                _cell_n = _tbl.cell(_row_start, 1)
                _cell_n.merge(_tbl.cell(_row_start + 3, 1))
                _cell_n.text = ""
                _cell_n.vertical_anchor = MSO_ANCHOR.MIDDLE
                _cell_n.margin_top = Emu(30000); _cell_n.margin_bottom = Emu(30000)
                _cell_n.margin_left = Emu(60000); _cell_n.margin_right = Emu(30000)
                _tf_n = _cell_n.text_frame; _tf_n.word_wrap = True
                for _p in list(_tf_n.paragraphs)[1:]:
                    _p._p.getparent().remove(_p._p)
                _p1 = _tf_n.paragraphs[0]; _p1.alignment = PP_ALIGN.LEFT
                _r1 = _p1.add_run(); _r1.text = _name
                _r1.font.name = "Calibri"; _r1.font.size = Pt(11); _r1.font.bold = True
                _r1.font.color.rgb = _BLACK
                _p2 = _tf_n.add_paragraph(); _p2.alignment = PP_ALIGN.LEFT
                _r2 = _p2.add_run(); _r2.text = _role
                _r2.font.name = "Calibri"; _r2.font.size = Pt(9); _r2.font.italic = True
                _r2.font.color.rgb = _ROLE

                for _off, _cat in enumerate(_CATS):
                    _r_abs = _row_start + _off
                    _is_sigma = (_off == 0)
                    _is_itr   = (_off == 2)
                    _is_tehn  = (_off == 3)
                    _row_fill = _SIGMA_FILL if _is_itr else None
                    _row_align = PP_ALIGN.RIGHT if _is_tehn else PP_ALIGN.LEFT
                    _v_prev, _v_curr = _cats_vals[_off]
                    _d_txt, _d_color = _chisl_delta_text(_v_prev, _v_curr)
                    _chisl_set_simple(_tbl.cell(_r_abs, 2), _cat,
                                      size=9, bold=_is_sigma, italic=_is_tehn,
                                      align=_row_align, fill=_row_fill)
                    _chisl_set_simple(_tbl.cell(_r_abs, 3), _chisl_fmt_int(_v_prev),
                                      size=(10 if _is_sigma else 9),
                                      align=_row_align, bold=_is_sigma, fill=_row_fill)
                    _chisl_set_simple(_tbl.cell(_r_abs, 4), _chisl_fmt_int(_v_curr),
                                      size=(10 if _is_sigma else 9),
                                      align=_row_align, bold=_is_sigma, fill=_row_fill)
                    _chisl_set_simple(_tbl.cell(_r_abs, 5), _d_txt,
                                      size=(10 if _is_sigma else 9),
                                      align=_row_align, bold=True, fill=_row_fill, color=_d_color)

            # ИТОГО — только на последнем чанке
            if _include_total:
                _r_total = 1 + len(_rows_chunk) * 4
                _cell_t_num = _tbl.cell(_r_total, 0)
                _cell_t_num.merge(_tbl.cell(_r_total + 3, 0))
                _chisl_set_simple(_cell_t_num, "", bold=True, size=12, fill=_TOTAL_FILL)
                _cell_t_name = _tbl.cell(_r_total, 1)
                _cell_t_name.merge(_tbl.cell(_r_total + 3, 1))
                _chisl_set_simple(_cell_t_name, "ИТОГО", bold=True, size=12,
                                  align=PP_ALIGN.RIGHT, fill=_TOTAL_FILL)
                for _off, _cat in enumerate(_CATS):
                    _r_abs = _r_total + _off
                    _is_sigma = (_off == 0)
                    _is_tehn  = (_off == 3)
                    _row_align = PP_ALIGN.RIGHT if _is_tehn else PP_ALIGN.LEFT
                    _v_prev, _v_curr = _total_cats[_off]
                    _d_txt, _d_color = _chisl_delta_text(_v_prev, _v_curr)
                    _chisl_set_simple(_tbl.cell(_r_abs, 2), _cat,
                                      size=10, bold=True, italic=_is_tehn,
                                      align=_row_align, fill=_TOTAL_FILL)
                    _chisl_set_simple(_tbl.cell(_r_abs, 3), _chisl_fmt_int(_v_prev),
                                      size=(11 if _is_sigma else 10),
                                      align=_row_align, bold=True, fill=_TOTAL_FILL)
                    _chisl_set_simple(_tbl.cell(_r_abs, 4), _chisl_fmt_int(_v_curr),
                                      size=(11 if _is_sigma else 10),
                                      align=_row_align, bold=True, fill=_TOTAL_FILL)
                    _chisl_set_simple(_tbl.cell(_r_abs, 5), _d_txt,
                                      size=(11 if _is_sigma else 10),
                                      align=_row_align, bold=True, fill=_TOTAL_FILL, color=_d_color)

        # Пагинация: по 9 подрядчиков на слайд (правка пользователя 20.05.2026).
        # Когда подрядчиков мало — один слайд, ИТОГО на нём. Когда больше — две
        # страницы, ИТОГО только на последней. Заголовок один и тот же на всех
        # страницах («РЕСУРСЫ»), приписки «— продолжение» нет.
        _CHISL_PER_SLIDE = 9
        _chisl_pages = [_chisl_rows[i:i + _CHISL_PER_SLIDE]
                        for i in range(0, len(_chisl_rows), _CHISL_PER_SLIDE)]
        if not _chisl_pages:
            _chisl_pages = [[]]  # пустой — на одну страницу
        _lc_idx = list(prs.slides).index(slide2)
        for _pg_idx, _chunk in enumerate(_chisl_pages):
            _slide = clone_slide(prs, slide2)
            if CHISL_SLIDE is None:
                CHISL_SLIDE = _slide
            for _sh in list(_slide.shapes):
                if getattr(_sh, 'name', '') not in _HEADER_NAMES_CHISL:
                    remove_shape(_sh)
            for _sh in _slide.shapes:
                if getattr(_sh, 'name', '') == 'Прямоугольник 11':
                    set_text(_sh, "РЕСУРСЫ"); break
            _is_last = (_pg_idx == len(_chisl_pages) - 1)
            _build_chisl_table(_slide, _chunk, _is_last)
            reorder_slide_to(prs, _slide, _lc_idx + 1 + _pg_idx)
        print(f"Ресурсы: {len(_chisl_rows)} подрядчиков → {len(_chisl_pages)} слайд(а)")
    except Exception as _e:
        print(f"WARN: Численность slide skipped: {_e}")
else:
    print("INFO: include_chislennost не задан — слайд Численности пропущен")

# ---------------------------- pagination + table ----------------------------
def update_title(slide, new_title):
    for s in slide.shapes:
        if s.has_text_frame and ('«РЕПИН»' in s.text_frame.text or 'РЕПИН' in s.text_frame.text or 'СКК' in s.text_frame.text):
            set_text(s, new_title); return

def _line_cost(item, name_cpl):
    # build_data_table обрезает имя до name_cpl символов и переноса нет —
    # каждая строка занимает ровно одну UNIFORM_ROW_H.
    return 1

def _section_key(it, show_contractor=True):
    """Ключ группировки для section-row. Объединяет (section, contractor):
    одна и та же секция работ с разными подрядчиками превращается в две
    группы (две section-row подряд). Правка 2026-05-20 для слайдов стадий —
    подтянуть подрядчика из МСГ (X на 'Раб. ГПР' строке).

    show_contractor=False (правка 2026-05-21 для слайда РД) — группируем
    только по section, подрядчика не учитываем. Это позволяет сливать в одну
    группу строки одного раздела с разными (или неуказанными) подрядчиками.
    """
    sec = it.get('section') or '— БЕЗ РАЗДЕЛА —'
    contr = (it.get('contractor') or '') if show_contractor else ''
    return (sec, contr)

def _page_weight(page, name_cpl=50, show_contractor=True):
    """Сколько строк (с учётом многострочного name) занимает page."""
    rows = 1  # header
    last_key = None
    for it in page:
        k = _section_key(it, show_contractor)
        if k != last_key:
            rows += 1; last_key = k
        rows += _line_cost(it, name_cpl)
    return rows

def paginate_by_section(items, max_rows=22, name_cpl=50, show_contractor=True):
    if not items: return []
    pages, cur, cur_rows, last_key = [], [], 1, None
    for it in items:
        k = _section_key(it, show_contractor)
        new_section = k != last_key
        cost = _line_cost(it, name_cpl)
        rows_needed = (1 if new_section else 0) + cost
        if cur_rows + rows_needed > max_rows and cur:
            pages.append(cur)
            cur, cur_rows, last_key = [], 1, None
            new_section = True; rows_needed = 1 + cost
        if new_section:
            last_key = k
        cur.append(it); cur_rows += rows_needed
    if cur: pages.append(cur)
    # Post-merge: если хвостовая страница вместе с предыдущей помещается
    # в max_rows — сливаем, чтобы не висел «одинокий» лист с 2-3 строками.
    while len(pages) >= 2 and _page_weight(pages[-2] + pages[-1], name_cpl, show_contractor) <= max_rows:
        pages[-2].extend(pages[-1])
        pages.pop()
    return pages

def build_data_table(slide, items, columns, top_emu=800000, name_cpl=50, stretch_to_fill=False, show_contractor=True):
    for s in list(slide.shapes):
        if s.shape_type == 13 and s.name != 'Рисунок 6':
            remove_shape(s)
        elif s.has_table:
            remove_shape(s)
    sw, sh = prs.slide_width, prs.slide_height
    L, T = Emu(341998), Emu(top_emu)
    W = Emu(sw - L*2); H = sh - T - Emu(300000)
    if not items:
        tb = slide.shapes.add_textbox(L, Emu(2942274), W, Emu(369332))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "ОТСТАВАНИЯ ОТ ПЛАНА ОТСУТСТВУЮТ"
        r.font.bold = True; r.font.size = Pt(24); r.font.name = 'Calibri'
        r.font.color.rgb = RGBColor(0xE6, 0x00, 0x73)
        return
    grouped = []
    for it in items:
        k = _section_key(it, show_contractor)
        if grouped and grouped[-1][0] == k:
            grouped[-1][1].append(it)
        else:
            grouped.append((k, [it]))
    total_rows = 1 + sum(1 + len(grp) for _, grp in grouped)
    # Если stretch_to_fill=True — растягиваем все строки страницы равномерно,
    # чтобы таблица заполнила лист до конца (нет белого разрыва снизу).
    # Высота строки всё равно одинаковая на странице, просто больше UNIFORM_ROW_H.
    if stretch_to_fill:
        row_h = max(UNIFORM_ROW_H, int((H - HEADER_ROW_H) / max(1, total_rows - 1)))
    else:
        row_h = UNIFORM_ROW_H
    table_h_emu = HEADER_ROW_H + (total_rows - 1) * row_h
    table_shape = slide.shapes.add_table(rows=total_rows, cols=len(columns), left=L, top=T, width=W, height=table_h_emu)
    tbl = table_shape.table
    props = [c[2] for c in columns]; sp = sum(props); props = [p/sp for p in props]
    for ci, p in enumerate(props):
        tbl.columns[ci].width = int(W * p)
    tbl.rows[0].height = HEADER_ROW_H
    for r in list(tbl.rows)[1:]:
        r.height = row_h
    for r in tbl.rows:
        for c in r.cells:
            c.margin_top = Emu(15000); c.margin_bottom = Emu(15000)
            c.margin_left = Emu(60000); c.margin_right = Emu(60000)
            c.text_frame.word_wrap = True
    for ci, (hdr, key, _, fmt) in enumerate(columns):
        cell = tbl.cell(0, ci); cell.text = ""
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(11); r.font.name = 'Calibri'
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    ri = 0
    for (sec, contr), grp in grouped:
        ri += 1
        first = tbl.cell(ri, 0); last = tbl.cell(ri, len(columns)-1)
        first.merge(last); first.text = ""
        # Правка пользователя 2026-05-20: подрядчик из 'Раб. ГПР' строки МСГ
        # подписан рядом с разделом через em-dash (если назначен). Имя
        # секции — navy bold, подрядчик — обычным шрифтом, тем же цветом.
        p = first.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = sec
        r.font.bold = True; r.font.size = Pt(11); r.font.name = 'Calibri'
        r.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        if contr:
            r2 = p.add_run(); r2.text = f"  —  {contr}"
            r2.font.bold = False; r2.font.size = Pt(11); r2.font.name = 'Calibri'
            r2.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        first.fill.solid(); first.fill.fore_color.rgb = RGBColor(0xDC, 0xE6, 0xF1)
        first.vertical_anchor = MSO_ANCHOR.MIDDLE
        for item in grp:
            ri += 1
            for ci, (hdr, key, _, fmt) in enumerate(columns):
                cell = tbl.cell(ri, ci); cell.text = ""
                p = cell.text_frame.paragraphs[0]
                v = item.get(key)
                txt = fmt(v) if fmt else (str(v) if v not in (None,'') else '–')
                if key == 'name' and len(txt) > name_cpl:
                    txt = txt[:max(1, name_cpl - 1)].rstrip() + '…'
                p.alignment = PP_ALIGN.LEFT if key == 'name' else PP_ALIGN.CENTER
                rn = p.add_run(); rn.text = txt
                rn.font.size = Pt(9); rn.font.name = 'Calibri'
                if key == 'days_dev' and isinstance(v, (int, float)) and v < 0:
                    rn.font.bold = True; rn.font.color.rgb = RGBColor(0xE0,0,0)
                if key == 'obj_otkl':
                    nv = _to_float_safe(v)
                    if nv is not None and nv < 0:
                        rn.font.bold = True; rn.font.color.rgb = RGBColor(0xE0,0,0)
                if key == 'prereq' and isinstance(v, str) and v.strip() == '-':
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xD9, 0xE1)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def build_paged_table(prs, base_slide, items, columns, top_emu, base_title, title_setter, max_rows=22, subtitle_fn=None, name_cpl=None, stretch_to_fill=False, show_contractor=True):
    # Если не передан name_cpl — берём из ширины первой колонки (это всегда «Наименование работ»):
    # 0.26 → ~40 символов на строку (СМР, 9 колонок), 0.32 → ~50 (не-СМР, 8 колонок).
    if name_cpl is None:
        name_w = next((c[2] for c in columns if c[1] == 'name'), 0.3)
        # k=0.0056 — эмпирический коэффициент для Calibri Pt(9) c учётом cell padding;
        # подбирался под максимальную плотность без вылезания «…» за правый край ячейки.
        name_cpl = max(28, int(name_w / 0.0056))
    pages = paginate_by_section(items, max_rows=max_rows, name_cpl=name_cpl, show_contractor=show_contractor)
    if not pages:
        # No lag → drop the slide entirely.
        remove_slide(prs, base_slide)
        return []
    base_idx = list(prs.slides).index(base_slide)
    slides = []
    for pi, page_items in enumerate(pages):
        if pi == 0:
            target = base_slide
        else:
            target = clone_slide(prs, base_slide)
            reorder_slide_to(prs, target, base_idx + pi)
            # 20.05.2026: суффикс «— продолжение» убран по запросу пользователя.
            # Continuation-слайд несёт идентичный заголовок (СМР/РД/...) — что
            # это вторая страница той же стадии, читателю и так очевидно.
            title_setter(target, base_title)
        if subtitle_fn:
            subtitle_fn(target)
        build_data_table(target, page_items, columns, top_emu=top_emu, name_cpl=name_cpl, stretch_to_fill=stretch_to_fill, show_contractor=show_contractor)
        slides.append(target)
    return slides

# (slide2..slide8 refs were captured at the top, before the ГПР insert.)

# ---------------------------- slide 3: СМР ----------------------------
update_title(slide3, f"СТРОИТЕЛЬНО МОНТАЖНЫЕ РАБОТЫ ({PROJECT_SHORT})")

def _set_title_smr(slide, text):
    for s in slide.shapes:
        if s.has_text_frame:
            t = s.text_frame.text.upper()
            if 'СМР' in t or 'СТРОИТЕЛЬНО' in t or 'МОНТАЖНЫЕ' in t:
                set_text(s, text); return

_SMR_SLIDES = build_paged_table(prs, slide3, SMR_LAG, [
    # Правка пользователя 20.05.2026: «Дни откл.» переехали ПОСЛЕ «Проект.
    # объём» — блок «объёмы» (Объём откл. / Ед. / Проект.объём) идёт цельно,
    # сразу за ним «Дни откл.», далее даты план/факт.
    ("Предш.", "prereq", 0.04, lambda v: str(v) if v else '–'),
    ("Наименование работ", "name", 0.46, lambda v: str(v)[:200]),
    ("Здан.", "block", 0.05, lambda v: str(v) if v else '–'),
    ("Объём\nоткл.", "obj_otkl", 0.07, fmt_num),
    ("Ед.", "unit", 0.05, lambda v: str(v) if v else '–'),
    ("Проект.\nобъём", "proj_obj", 0.07, fmt_pos),
    ("Дни\nоткл.", "days_dev", 0.05, fmt_days_neg_only),
    ("План\nначало", "plan_n", 0.07, fmt_date),
    ("План\nоконч.", "plan_o", 0.07, fmt_date),
    ("Факт\nначало", "fact_n", 0.07, fmt_date),
], top_emu=800000, base_title=f"СТРОИТЕЛЬНО МОНТАЖНЫЕ РАБОТЫ ({PROJECT_SHORT})", title_setter=_set_title_smr, max_rows=30)

# ---------------------------- slide 3.5: Монолит (только для Репино) ----------------------------
# Логика 1-в-1 с /мсг (build_meeting_report.py): 13 шаблонных строк × 6 корпусов × {План, Факт},
# матчинг по name/section/subsection, фильтр по текущему месяцу. Источник — items.json.
# Вставляется сразу ПОСЛЕ последнего СМР-слайда, ПЕРЕД РД.
MONOLIT_SLIDE = None
if CFG.get("include_monolit") and _SMR_SLIDES:
    _RU_MONTHS = {1:"Январь",2:"Февраль",3:"Март",4:"Апрель",5:"Май",6:"Июнь",
                  7:"Июль",8:"Август",9:"Сентябрь",10:"Октябрь",11:"Ноябрь",12:"Декабрь"}
    CURRENT_MONTH = _RU_MONTHS[TODAY.month]
    BUILDINGS = ("К1","К2","К3","К4","К5","К6")

    def _has(text, *needles):
        if not isinstance(text, str): return False
        s = text.lower()
        return all((n.lower() in s) for n in needles)

    # Обновлено по файлу «Парапет.pptx» (правка ОЛ 19.05.2026):
    # — подписи -1 эт. → «(подвал)»;
    # — этажность перекрытий сдвинута на −1 («2 эт.» → «1 эт.» и т.д.) —
    #   перекрытие на отметке N это потолок N-го этажа, а не пола следующего;
    # — добавлена 14-я строка «ПАРАПЕТ КРОВЛИ».
    # Логика матчинга по name/section_rd/subsection_rd НЕ менялась — это были
    # только подписи на слайде.
    # Правка пользователя 20.05.2026: перекрытия матчим по «ПП<N>» (индекс
    # плиты — единственный надёжный признак, отметки в gsheet могут не совпадать
    # с шаблонными: подвал в gsheet «ПП0 на отм. -0.200», в шаблоне «-0.100»).
    # В label через запятую добавлен «ПП<N>», чтобы было видно по чему матчим.
    MONOLIT_TEMPLATE = [
        ("ФУНДАМЕНТНАЯ ПЛИТА",
         lambda x, d, e: _has(x, "фундамент")),
        ("МОНОЛИТНЫЕ Ж/Б СТЕНЫ НИЖЕ ОТМ. 0.000 (подвал)",
         lambda x, d, e: ("Монолит ниже 0" in (d or "")) and ("Монолитные стены" in (e or ""))),
        ("МОНОЛИТНЫЕ Ж/Б ПЕРЕКРЫТИЯ НА ОТМ. -0.100 (подвал, ПП0)",
         lambda x, d, e: _has(x, "перекрыти", "ПП0")),
        ("МОНОЛИТНЫЕ Ж/Б СТЕНЫ НИЖЕ ОТМ. +3.200 (1 эт.)",
         lambda x, d, e: _has(x, "стен", "1-ого этажа")),
        ("МОНОЛИТНЫЕ Ж/Б ПЕРЕКРЫТИЯ НА ОТМ. +3.200 (1 эт., ПП1)",
         lambda x, d, e: _has(x, "перекрыти", "ПП1")),
        ("МОНОЛИТНЫЕ Ж/Б СТЕНЫ НИЖЕ ОТМ. +6.150 (2 эт.)",
         lambda x, d, e: _has(x, "стен", "2-ого этажа")),
        ("МОНОЛИТНЫЕ Ж/Б ПЕРЕКРЫТИЯ НА ОТМ. +6.150 (2 эт., ПП2)",
         lambda x, d, e: _has(x, "перекрыти", "ПП2")),
        ("МОНОЛИТНЫЕ Ж/Б СТЕНЫ НИЖЕ ОТМ. +9.800 (3 эт.)",
         lambda x, d, e: _has(x, "стен", "3-ого этажа")),
        ("МОНОЛИТНЫЕ Ж/Б ПЕРЕКРЫТИЯ НА ОТМ. +9.800 (3 эт., ПП3)",
         lambda x, d, e: _has(x, "перекрыти", "ПП3")),
        ("МОНОЛИТНЫЕ Ж/Б СТЕНЫ НИЖЕ ОТМ. +13.400 (4 эт.)",
         lambda x, d, e: _has(x, "стен", "4-ого этажа")),
        ("МОНОЛИТНЫЕ Ж/Б ПЕРЕКРЫТИЯ НА ОТМ. +13.400 (4 эт., ПП4)",
         lambda x, d, e: _has(x, "перекрыти", "ПП4")),
        ("МОНОЛИТНЫЕ Ж/Б СТЕНЫ НИЖЕ ОТМ. +16.700 (5 эт.)",
         lambda x, d, e: _has(x, "стен", "5-ого этажа")),
        ("МОНОЛИТНЫЕ Ж/Б ПЕРЕКРЫТИЯ НА ОТМ. +16.700 (5 эт., ПП5)",
         lambda x, d, e: _has(x, "перекрыти", "ПП5")),
        ("ПАРАПЕТ КРОВЛИ",
         lambda x, d, e: _has(x, "парапет")),
    ]

    # ОЛ 19.05.2026: Откл берём не как «Факт − План», а из колонки AA МСГ
    # «Объём откл./опереж. на сег.» (это `obj_otkl` в items.json) — суммарно
    # по всем строкам 'План' данного матч-блока. Отриц. значение = отстаём.
    _mtable = [{b: {"План": 0.0, "Факт": 0.0, "Откл": 0.0,
                    "ОтклЕсть": False} for b in BUILDINGS}
               for _ in MONOLIT_TEMPLATE]
    # Правка пользователя 20.05.2026: фильтр по month_filter (колонка O МСГ)
    # снят. AQ («План/Факт МЕСЯЦ») в gsheet — формула, которая сама
    # агрегирует значения текущего месяца из календаря; если строка не активна
    # в этом месяце, AQ пуст/0 и до таблицы не доходит. Колонка O — рабочий
    # фильтр planner'а, не всегда заполнена; без неё мы корректно ловим всё
    # что реально идёт в текущем месяце.
    _matched = 0
    for it in ITEMS:
        x  = it.get('name')
        y  = it.get('block')
        d  = it.get('section_rd')
        e  = it.get('subsection_rd')
        al = it.get('rtype')
        aq = it.get('pf_month')
        if not isinstance(x, str) or "етонирование" not in x: continue
        if y not in BUILDINGS: continue
        if al not in ("План", "Факт"): continue
        # pf_month — План/Факт месяца (для накопления «План м³» / «Факт м³»)
        try:
            v = _to_float(aq) if aq not in (None, "") else None
        except (TypeError, ValueError):
            v = None
        for _idx, (_lbl, matcher) in enumerate(MONOLIT_TEMPLATE):
            try:
                if matcher(x, d, e):
                    if v is not None:
                        _mtable[_idx][y][al] += v
                        _matched += 1
                    if al == 'План':
                        # obj_otkl — отклонение НА СЕГ. (только в План-строках)
                        _oo = _to_float_safe(it.get('obj_otkl'))
                        if _oo is not None:
                            _mtable[_idx][y]["Откл"] += _oo
                            _mtable[_idx][y]["ОтклЕсть"] = True
                    break
            except Exception:
                pass
    print(f"Монолит: matched {_matched} строк под {CURRENT_MONTH}")

    # Клонируем последний СМР-слайд, оставляем только хедерные шейпы — табличку отрисуем сами.
    _last_smr = _SMR_SLIDES[-1]
    MONOLIT_SLIDE = clone_slide(prs, _last_smr)
    _HEADER_NAMES_MON = {
        'Прямая соединительная линия 10',
        'Прямоугольник 11',
        'Прямоугольник 12',
        'Рисунок 6',
    }
    for _sh in list(MONOLIT_SLIDE.shapes):
        if getattr(_sh, 'name', '') not in _HEADER_NAMES_MON:
            remove_shape(_sh)
    # Заголовок
    for _sh in MONOLIT_SLIDE.shapes:
        if getattr(_sh, 'name', '') == 'Прямоугольник 11':
            set_text(_sh, f"МОНОЛИТ ({PROJECT_SHORT})")
            break

    # ОЛ 19.05.2026: добавлена колонка «Откл.» в каждую тройку План/Факт/Откл
    # на корпус — итого 20 колонок (B,C + 6 корпусов × 3). Зелёная подсветка
    # при Факт=План (выполнено), красная при Факт<План (отставание).
    # Таблица: 3 строки шапки + 13 строк данных + 1 ИТОГО = 17 строк.
    _n_tpl = len(MONOLIT_TEMPLATE)
    _n_rows = 3 + _n_tpl + 1
    _n_cols = 2 + 6 * 3  # 20
    _sw_e, _sh_e = prs.slide_width, prs.slide_height
    _L = Emu(341998); _T = Emu(900000)
    _W = _sw_e - _L * 2
    _H = _sh_e - _T - Emu(250000)

    _tbl_shape = MONOLIT_SLIDE.shapes.add_table(rows=_n_rows, cols=_n_cols, left=_L, top=_T, width=_W, height=_H)
    _tbl = _tbl_shape.table

    # Ширины: B (Вид) 0.30, C (Месяц) 0.05, далее 18 одинаковых План/Факт/Откл колонок.
    _BC_W = int(_W * 0.30); _C_W = int(_W * 0.05)
    _PF_W = int((_W - _BC_W - _C_W) / 18)
    _tbl.columns[0].width = _BC_W
    _tbl.columns[1].width = _C_W
    for _ci in range(2, _n_cols):
        _tbl.columns[_ci].width = _PF_W

    # Высоты строк: шапка ~250K EMU, данные плотные ~140K (~14.6 pt)
    _HDR_H = 250000
    _DATA_H = 140000
    for _r in range(3):
        _tbl.rows[_r].height = Emu(_HDR_H)
    for _r in range(3, _n_rows):
        _tbl.rows[_r].height = Emu(_DATA_H)

    _HDR_RGB     = RGBColor(0x1F, 0x4E, 0x79)  # navy — шапка таблицы (3 верхние строки)
    _PLAN_RGB    = RGBColor(0xDC, 0xE6, 0xF1)  # голубой — План-колонки данных и ИТОГО
    _FACT_RGB    = RGBColor(0xEA, 0xF1, 0xF8)  # светло-голубой — Факт/Откл-колонки, B/C данных
    _DL_RED      = RGBColor(0xC0, 0x00, 0x00)  # цвет ШРИФТА Откл при отставании (правка ОЛ 19.05.2026)
    _BLACK_RGB   = RGBColor(0x00, 0x00, 0x00)
    _WHITE_RGB   = RGBColor(0xFF, 0xFF, 0xFF)
    _BORDER_HEX  = "BFBFBF"   # светло-серый — тонкие ровные границы
    _BORDER_W    = 6350       # 0.5 pt в EMU

    def _set_cell_border(cell, color_hex=_BORDER_HEX, width_emu=_BORDER_W):
        """Тонкие рамки на всех 4 сторонах ячейки (python-pptx нативно не умеет — пишем в OOXML).

        Важно: по схеме CT_TableCellProperties элементы lnL/lnR/lnT/lnB должны идти ПЕРЕД
        fill-элементами (solidFill и т.п.). SubElement добавляет в хвост — если fill уже
        стоит, схема ломается и PowerPoint молча игнорирует рамки (отсюда «белые рамки»
        в первой попытке). Чиним через insert(0, ...) в обратном порядке.
        """
        tcPr = cell._tc.get_or_add_tcPr()
        for _side in ('lnL', 'lnR', 'lnT', 'lnB'):
            for _existing in tcPr.findall(qn(f'a:{_side}')):
                tcPr.remove(_existing)
        # Вставляем в обратном порядке через insert(0,...) — финальный порядок: lnL, lnR, lnT, lnB
        for _side in ('lnB', 'lnT', 'lnR', 'lnL'):
            _ln = etree.Element(qn(f'a:{_side}'))
            _ln.set('w', str(width_emu))
            _ln.set('cap', 'flat')
            _ln.set('cmpd', 'sng')
            _ln.set('algn', 'ctr')
            _sf = etree.SubElement(_ln, qn('a:solidFill'))
            _sc = etree.SubElement(_sf, qn('a:srgbClr')); _sc.set('val', color_hex)
            _pd = etree.SubElement(_ln, qn('a:prstDash')); _pd.set('val', 'solid')
            etree.SubElement(_ln, qn('a:round'))
            tcPr.insert(0, _ln)

    def _set_cell(cell, text, *, bold=False, size=9, align=PP_ALIGN.CENTER, fill=None, color=None):
        # Auto-contrast: на navy-шапке — белый текст, на голубых — чёрный.
        if color is None:
            color = _WHITE_RGB if (fill is not None and fill == _HDR_RGB) else _BLACK_RGB
        cell.text = ""
        p = cell.text_frame.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(text)
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        cell.text_frame.word_wrap = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Emu(2000); cell.margin_bottom = Emu(2000)
        cell.margin_left = Emu(20000); cell.margin_right = Emu(20000)
        if fill is not None:
            cell.fill.solid(); cell.fill.fore_color.rgb = fill

    # Row 0 — «Вид монолита» / «Месяц» / «№ корпуса» — шапка navy
    _set_cell(_tbl.cell(0, 0), "Вид монолита", bold=True, size=11, fill=_HDR_RGB)
    _tbl.cell(0, 0).merge(_tbl.cell(2, 0))   # vertical merge через 3 строки шапки
    _set_cell(_tbl.cell(0, 1), "Месяц",       bold=True, size=11, fill=_HDR_RGB)
    _tbl.cell(0, 1).merge(_tbl.cell(2, 1))
    _set_cell(_tbl.cell(0, 2), "№ корпуса",   bold=True, size=11, fill=_HDR_RGB)
    _tbl.cell(0, 2).merge(_tbl.cell(0, _n_cols - 1))  # horizontal merge через все корпусные колонки

    # Row 1 — К1..К6 (с приставкой «К» — правка ОЛ 19.05.2026)
    for _i in range(6):
        _c = 2 + _i * 3
        _set_cell(_tbl.cell(1, _c), f"К{_i + 1}", bold=True, size=11, fill=_HDR_RGB)
        _tbl.cell(1, _c).merge(_tbl.cell(1, _c + 2))  # merge всех 3 колонок тройки

    # Row 2 — «План, м³» / «Факт, м³» / «Откл.»
    for _i in range(6):
        _c = 2 + _i * 3
        _set_cell(_tbl.cell(2, _c),     "План, м³", bold=True, size=10, fill=_PLAN_RGB)
        _set_cell(_tbl.cell(2, _c + 1), "Факт, м³", bold=True, size=10, fill=_FACT_RGB)
        _set_cell(_tbl.cell(2, _c + 2), "Откл.",    bold=True, size=10, fill=_FACT_RGB)

    def _fmt_int_signed(v):
        """Целое со знаком. ОЛ 19.05.2026: положительные значения (опережение
        плана) на слайде не показываем — пусто. Показываем только 0 и
        отрицательные (отставание)."""
        if v is None: return ""
        n = int(round(v))
        if n == 0: return "0"
        if n > 0: return ""
        return f"{n:+d}"

    def _fmt_int_plain(v):
        if v is None or abs(v) < 0.5: return ""
        return str(int(round(v)))

    # ОЛ 19.05.2026: заливку «зелёный/красный» убираем — фон везде нейтральный
    # (План: _PLAN_RGB, Факт/Откл: _FACT_RGB). Подсветка отставания — только
    # шрифт Откл-колонки красным при отрицательном значении.
    # Rows 3..15 — данные шаблона
    for _ti, (_label, _matcher) in enumerate(MONOLIT_TEMPLATE):
        _r = 3 + _ti
        _set_cell(_tbl.cell(_r, 0), _label, bold=False, size=8, align=PP_ALIGN.LEFT, fill=_FACT_RGB)
        _set_cell(_tbl.cell(_r, 1), CURRENT_MONTH, bold=False, size=9, fill=_FACT_RGB)
        for _bi, _b in enumerate(BUILDINGS):
            _c = 2 + _bi * 3
            _v_pl = _mtable[_ti][_b]["План"]
            _v_fa = _mtable[_ti][_b]["Факт"]
            _v_dl = _mtable[_ti][_b]["Откл"] if _mtable[_ti][_b]["ОтклЕсть"] else None
            _txt_pl = _fmt_int_plain(_v_pl)
            _txt_fa = _fmt_int_plain(_v_fa)
            _txt_dl = _fmt_int_signed(_v_dl)
            _dl_color = _DL_RED if (_v_dl is not None and _v_dl < -0.5) else None
            _set_cell(_tbl.cell(_r, _c),     _txt_pl, size=9, align=PP_ALIGN.CENTER, fill=_PLAN_RGB)
            _set_cell(_tbl.cell(_r, _c + 1), _txt_fa, size=9, align=PP_ALIGN.CENTER, fill=_FACT_RGB)
            _set_cell(_tbl.cell(_r, _c + 2), _txt_dl, size=9,
                      bold=(_dl_color is not None),
                      align=PP_ALIGN.CENTER, fill=_FACT_RGB, color=_dl_color)

    # Толстый вертикальный разделитель между корпусами (правка ОЛ 19.05.2026).
    # Идёт по правой границе колонок 4, 7, 10, 13, 16 (последняя «Откл» каждого
    # из К1..К5) — отделяет тройку План/Факт/Откл одного корпуса от другого.
    _COL_SEP_HEX   = "1F4E79"   # navy
    _COL_SEP_W_EMU = 25400      # 2.0 pt в EMU

    def _set_cell_border_right_thick(cell, color_hex, width_emu):
        tcPr = cell._tc.get_or_add_tcPr()
        for _existing in tcPr.findall(qn('a:lnR')):
            tcPr.remove(_existing)
        _ln = etree.Element(qn('a:lnR'))
        _ln.set('w', str(width_emu))
        _ln.set('cap', 'flat')
        _ln.set('cmpd', 'sng')
        _ln.set('algn', 'ctr')
        _sf = etree.SubElement(_ln, qn('a:solidFill'))
        _sc = etree.SubElement(_sf, qn('a:srgbClr')); _sc.set('val', color_hex)
        _pd = etree.SubElement(_ln, qn('a:prstDash')); _pd.set('val', 'solid')
        etree.SubElement(_ln, qn('a:round'))
        tcPr.insert(0, _ln)

    # ИТОГО — суммы по каждой тройке (нейтральный фон, шрифт Откл красным при отриц.)
    _r_total = 3 + _n_tpl
    _set_cell(_tbl.cell(_r_total, 0), "ИТОГО:", bold=True, size=10, align=PP_ALIGN.RIGHT, fill=_FACT_RGB)
    _set_cell(_tbl.cell(_r_total, 1), "",       bold=True, size=10, fill=_FACT_RGB)
    for _bi, _b in enumerate(BUILDINGS):
        _c = 2 + _bi * 3
        _sum_pl = sum(_mtable[_ti][_b]["План"] for _ti in range(_n_tpl))
        _sum_fa = sum(_mtable[_ti][_b]["Факт"] for _ti in range(_n_tpl))
        _has_dl_any = any(_mtable[_ti][_b]["ОтклЕсть"] for _ti in range(_n_tpl))
        _sum_dl = sum(_mtable[_ti][_b]["Откл"] for _ti in range(_n_tpl)) if _has_dl_any else None
        _txt_pl = _fmt_int_plain(_sum_pl)
        _txt_fa = _fmt_int_plain(_sum_fa)
        _txt_dl = _fmt_int_signed(_sum_dl)
        _dl_color = _DL_RED if (_sum_dl is not None and _sum_dl < -0.5) else None
        _set_cell(_tbl.cell(_r_total, _c),     _txt_pl, bold=True, size=10, align=PP_ALIGN.CENTER, fill=_PLAN_RGB)
        _set_cell(_tbl.cell(_r_total, _c + 1), _txt_fa, bold=True, size=10, align=PP_ALIGN.CENTER, fill=_FACT_RGB)
        _set_cell(_tbl.cell(_r_total, _c + 2), _txt_dl, bold=True, size=10, align=PP_ALIGN.CENTER, fill=_FACT_RGB, color=_dl_color)

    # Толстая правая граница на колонках 4, 7, 10, 13, 16 — это Откл-колонка
    # каждого К1..К5, после неё начинается следующий корпус. Применяем ко
    # всем строкам таблицы (шапка + данные + ИТОГО).
    _SEPARATOR_COLS = [4, 7, 10, 13, 16]
    for _row_idx in range(_n_rows):
        for _sep_col in _SEPARATOR_COLS:
            _set_cell_border_right_thick(_tbl.cell(_row_idx, _sep_col),
                                         _COL_SEP_HEX, _COL_SEP_W_EMU)

    # Позиционируем Монолит-слайд сразу после последнего СМР-слайда
    _last_smr_idx = list(prs.slides).index(_last_smr)
    reorder_slide_to(prs, MONOLIT_SLIDE, _last_smr_idx + 1)
    print(f"Монолит-слайд вставлен после СМР (idx {_last_smr_idx + 1})")

# ---------------------------- slide 4: РД ----------------------------
update_title(slide4, f"РАБОЧАЯ ДОКУМЕНТАЦИЯ ({PROJECT_SHORT})")
for s in list(slide4.shapes):
    if s.has_text_frame and ('ПОЛНОМ' in s.text_frame.text or 'ОТСТАЁТ' in s.text_frame.text):
        remove_shape(s)

def _set_title_rd(slide, text):
    for s in slide.shapes:
        if s.has_text_frame and ('РАБОЧАЯ' in s.text_frame.text.upper() or 'РАБОЧАЙ' in s.text_frame.text.upper()):
            set_text(s, text); return

# ОЛ 19.05.2026: «Дни откл.» переехали из хвоста сразу за «Ед.» —
# чтобы Объём + Ед. + Дни читались одним блоком. Применено к РД, Пакету,
# Тендеру, Договору, Мобилизации (и СМР выше — там собственный список).
_LIFECYCLE_COLS = [
    ("Предш.", "prereq", 0.04, lambda v: str(v) if v else '–'),
    ("Наименование работ", "name", 0.35, lambda v: str(v)[:220]),
    ("Здан.", "block", 0.07, lambda v: str(v) if v else '–'),
    ("Объём\nоткл.", "obj_otkl", 0.10, fmt_num),
    ("Ед.", "unit", 0.07, lambda v: str(v) if v else '–'),
    ("Дни\nоткл.", "days_dev", 0.07, fmt_days_neg_only),
    ("План\nначало", "plan_n", 0.10, fmt_date),
    ("План\nоконч.", "plan_o", 0.10, fmt_date),
    ("Факт\nначало", "fact_n", 0.10, fmt_date),
]

# РД — без колонки «Предш.»: в gsheet колонка V на стадии РД всегда пуста
# (РД — первая стадия цикла, у неё нет предшественника).
_RD_COLS = [
    ("Наименование работ", "name", 0.39, lambda v: str(v)[:220]),
    ("Здан.", "block", 0.07, lambda v: str(v) if v else '–'),
    ("Объём\nоткл.", "obj_otkl", 0.10, fmt_num),
    ("Ед.", "unit", 0.07, lambda v: str(v) if v else '–'),
    ("Дни\nоткл.", "days_dev", 0.07, fmt_days_neg_only),
    ("План\nначало", "plan_n", 0.10, fmt_date),
    ("План\nоконч.", "plan_o", 0.10, fmt_date),
    ("Факт\nначало", "fact_n", 0.10, fmt_date),
]

build_paged_table(prs, slide4, RD_LAG, _RD_COLS,
    top_emu=800000, base_title=f"РАБОЧАЯ ДОКУМЕНТАЦИЯ ({PROJECT_SHORT})",
    title_setter=_set_title_rd, max_rows=30,
    # Правка пользователя 2026-05-21: на РД-слайде подрядчика на section-row
    # не показывать. На стадии РД подрядчик ещё не релевантен (это будущий
    # генподрядчик СМР, не исполнитель РД), и section-rows одного раздела
    # перестают дублироваться из-за разных «будущих» подрядчиков.
    show_contractor=False)

# ---------------------------- slide 4b: ПАКЕТ ----------------------------
# Only create a Пакет slide if there's lag.
if PAK_LAG:
    slide_pak = clone_slide(prs, slide4)
    rd_last_idx = list(prs.slides).index(slide4) if slide4 in list(prs.slides) else -1
    def _is_rd_slide(sl):
        for sh in sl.shapes:
            if sh.has_text_frame and 'РАБОЧАЯ ДОКУМЕНТАЦИЯ' in sh.text_frame.text.upper():
                return True
        return False
    if rd_last_idx >= 0:
        i = rd_last_idx + 1
        slides_list = list(prs.slides)
        while i < len(slides_list) and slides_list[i] is not slide_pak and _is_rd_slide(slides_list[i]):
            i += 1
        reorder_slide_to(prs, slide_pak, i)
    for s in list(slide_pak.shapes):
        if s.has_text_frame and ('ОТСТАЁТ' in s.text_frame.text or 'ОПЕРЕЖАЕТ' in s.text_frame.text):
            remove_shape(s)
    for s in slide_pak.shapes:
        if s.has_text_frame and 'РАБОЧАЯ' in s.text_frame.text.upper():
            set_text(s, f"ПАКЕТ ({PROJECT_SHORT})")
            break

    def _set_title_pak(slide, text):
        for s in slide.shapes:
            if s.has_text_frame and ('ПАКЕТА' in s.text_frame.text.upper() or 'ПАКЕТ' in s.text_frame.text.upper()):
                set_text(s, text); return

    build_paged_table(prs, slide_pak, PAK_LAG, _LIFECYCLE_COLS,
        top_emu=800000, base_title=f"ПАКЕТ ({PROJECT_SHORT})",
        title_setter=_set_title_pak, max_rows=30)

# ---------------------------- slide 5: Тендер ----------------------------
if TND_LAG:
    update_title(slide5, f"ТЕНДЕР ({PROJECT_SHORT})")
    for s in list(slide5.shapes):
        if s.has_text_frame and ('ИНФОРМАЦИЯ' in s.text_frame.text or 'ОТСТАВАНИЯ' in s.text_frame.text or 'ПОЗИЦИЙ' in s.text_frame.text):
            remove_shape(s)
    def _set_title_tnd(slide, text):
        for s in slide.shapes:
            if s.has_text_frame and 'ТЕНДЕР' in s.text_frame.text.upper():
                set_text(s, text); return
    build_paged_table(prs, slide5, TND_LAG, _LIFECYCLE_COLS,
        top_emu=800000, base_title=f"ТЕНДЕР ({PROJECT_SHORT})",
        title_setter=_set_title_tnd, max_rows=30)
else:
    remove_slide(prs, slide5)

# ---------------------------- slide 6: Договор ----------------------------
update_title(slide6, f"ДОГОВОР ({PROJECT_SHORT})")
for s in list(slide6.shapes):
    if hasattr(s, 'name') and s.name == 'Прямая соединительная линия 4':
        remove_shape(s)

def _set_title_dog(slide, text):
    for s in slide.shapes:
        if s.has_text_frame and 'ДОГОВОР' in s.text_frame.text.upper():
            set_text(s, text); return

build_paged_table(prs, slide6, DOG_LAG, _LIFECYCLE_COLS,
    top_emu=800000, base_title=f"ДОГОВОР ({PROJECT_SHORT})",
    title_setter=_set_title_dog, max_rows=30)

# ---------------------------- slide 7: Финансирование (ОТКЛЮЧЁН) ----------------------------
# Слайд «Финансирование» временно убран из колоды (2026-05-12).
# Шаблон template/repino_template.pptx по-прежнему содержит этот слайд седьмым,
# поэтому здесь его нужно удалить из деки. Полная логика по построению таблицы
# с данными по финансированию сохранена в build_report_v_with_finance.py.
remove_slide(prs, slide7)

# ---------------------------- slide 8: Мобилизация ----------------------------
update_title(slide8, f"МОБИЛИЗАЦИЯ ({PROJECT_SHORT})")
for s in list(slide8.shapes):
    if hasattr(s, 'name') and s.name == 'Прямая соединительная линия 4':
        remove_shape(s)

def _set_title_mob(slide, text):
    for s in slide.shapes:
        if s.has_text_frame and 'МОБИЛИЗАЦИЯ' in s.text_frame.text.upper():
            set_text(s, text); return

build_paged_table(prs, slide8, MOB_LAG, _LIFECYCLE_COLS,
    top_emu=800000, base_title=f"МОБИЛИЗАЦИЯ ({PROJECT_SHORT})",
    title_setter=_set_title_mob, max_rows=30)

# ---------------------------- разделительные слайды перед стадиями ----------------------------
# Перед каждой стадией (СМР, РД, Пакет, Тендер, Договор, Мобилизация) — пустой
# слайд с крупным названием по центру. Если у стадии не оказалось отставаний
# (base_slide уже в SLIDES_TO_REMOVE) — разделитель не создаём.
# Подложка-картинка прозрачностью 20% берётся из `МСГ\dividers\<slug>.png`
# (если файла нет — слайд без подложки).
#
# 20.05.2026: ОТКЛЮЧЕНО ПО ЗАПРОСУ ПОЛЬЗОВАТЕЛЯ — заголовок самой стадии уже
# даёт достаточный визуальный разрыв. Код оставлен в скрипте — включается
# флагом `include_stage_dividers: true` в `projects\<key>\config.json`.
if CFG.get("include_stage_dividers"):
    _slide_pak_ref = slide_pak if PAK_LAG else None
    _DIVIDER_SPECS = [
        # (title, base_slide, bg_slug)
        ("СТРОИТЕЛЬНО МОНТАЖНЫЕ РАБОТЫ",    slide3,         "smr"),
        # МОНОЛИТ — без отдельного divider (Монолит — часть СМР, divider избыточен).
        ("РАБОЧАЯ ДОКУМЕНТАЦИЯ",            slide4,         "rd"),
        ("ПАКЕТ",                           _slide_pak_ref, "paket"),
        ("ТЕНДЕР",                          slide5,         "tender"),
        ("ДОГОВОР",                         slide6,         "dogovor"),
        ("МОБИЛИЗАЦИЯ",                     slide8,         "mob"),
    ]
    _DIVIDERS_DIR = os.path.normpath(os.path.join(
        os.path.dirname(os.path.abspath(__file__)), '..', 'dividers'))

    def _add_divider_background(slide, image_path):
        """Полнослайдовая подложка с прозрачностью 20% (как на обложке). Картинка
        вставляется первой в spTree, чтобы быть позади всех остальных шейпов."""
        pic = slide.shapes.add_picture(image_path, Emu(0), Emu(0),
                                       prs.slide_width, prs.slide_height)
        bf = pic._element.find('.//' + qn('p:blipFill'))
        if bf is not None:
            blip = bf.find(qn('a:blip'))
            if blip is not None:
                a = etree.SubElement(blip, qn('a:alphaModFix'))
                a.set('amt', '20000')
        sp_tree = pic._element.getparent()
        sp_tree.remove(pic._element)
        insert_idx = 0
        for i, child in enumerate(sp_tree):
            tag = etree.QName(child).localname
            if tag in ('nvGrpSpPr', 'grpSpPr'):
                insert_idx = i + 1
        sp_tree.insert(insert_idx, pic._element)

    for _title, _base, _bg_slug in _DIVIDER_SPECS:
        if _base is None or _base in SLIDES_TO_REMOVE: continue
        _div = clone_slide(prs, _base)
        for _sh in list(_div.shapes):
            remove_shape(_sh)
        _bg_path = os.path.join(_DIVIDERS_DIR, f"{_bg_slug}.png")
        if os.path.exists(_bg_path):
            try:
                _add_divider_background(_div, _bg_path)
                print(f"Divider '{_title}': подложка {_bg_slug}.png подключена")
            except Exception as _e:
                print(f"WARN: Divider '{_title}': не удалось добавить подложку: {_e}")
        else:
            print(f"INFO: Divider '{_title}': файл {_bg_path} отсутствует — без подложки")
        _sw, _sh_h = prs.slide_width, prs.slide_height
        _tb = _div.shapes.add_textbox(Emu(0), _sh_h // 2 - Emu(700000), _sw, Emu(1400000))
        _tf = _tb.text_frame; _tf.word_wrap = True
        _p = _tf.paragraphs[0]; _p.alignment = PP_ALIGN.CENTER
        _r = _p.add_run(); _r.text = _title
        _r.font.bold = True; _r.font.size = Pt(54); _r.font.name = 'Calibri'
        _r.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        _base_idx = list(prs.slides).index(_base)
        reorder_slide_to(prs, _div, _base_idx)

# ---------------------------- Предписания summary slide (последний) ----------------------------
PRED_FILE_ID = CFG.get("predpisaniya_summary_file_id")
if PRED_FILE_ID:
    try:
        import tempfile as _tf
        _pred_src = os.path.join(_tf.gettempdir(), f"pred_src_{PROJECT_KEY}.xlsx")
        _bitrix_download_disk_file(PRED_FILE_ID, _pred_src)
        _pred_rows = _pred_load_filtered_rows(_pred_src)
        if _pred_rows:
            # ОЛ 19.05.2026 (правка пользователя): «Статус устранения» заменён
            # на «Отметка об устранении / прогнозируемая дата устранения»
            # (колонка K из исходного xlsx). «Статус снятия» рядом (placeholder
            # «–»). Шрифт везде Calibri 9pt — без 8pt для длинного текста.
            _PRED_HDR     = ["№\nпредп.", "Дата\nвыдачи", "Подрядчик", "Перечень нарушений",
                             "Срок\nустранения", "Срыв\n(дней)",
                             "Отметка об устранении /\nпрогнозируемая дата устр.",
                             "Статус\nснятия"]
            # Правка пользователя 20.05.2026: «Отметка об устранении» в 2 раза
            # уже (0.30 → 0.15), освободившаяся доля уходит в «Перечень
            # нарушений» (0.31 → 0.46).
            _PRED_FRACS   = [0.04, 0.07, 0.08, 0.46, 0.07, 0.05, 0.15, 0.08]
            # Правка пользователя 2026-05-20: «Отметка об устранении» (ci=6)
            # — выравнивание по центру (раньше LEFT). Содержимое — даты-фразы
            # типа «Устранено 03.04.2026», логично центрировать как «Срок» и
            # «Срыв».
            _PRED_ALIGNS  = [PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT,
                             PP_ALIGN.LEFT,  PP_ALIGN.CENTER, PP_ALIGN.CENTER,
                             PP_ALIGN.CENTER, PP_ALIGN.CENTER]
            _PRED_NAVY    = RGBColor(0x1F, 0x4E, 0x79)
            _PRED_BLACK   = RGBColor(0x00, 0x00, 0x00)
            _PRED_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
            _PRED_RED     = RGBColor(0xC0, 0x00, 0x00)
            # Подсветки правки 20.05.2026:
            # — «Статус снятия» = «Не снято» → светло-розовая заливка, чёрный текст;
            # — «Отметка об устранении» содержит подстроку «Устранено» (дословно,
            #   с заглавной У) → светло-зелёная заливка. Формы «Устранены»/
            #   «Устранение» — НЕ подсвечивать (по явному запросу).
            _PRED_PINK    = RGBColor(0xFF, 0xD9, 0xE1)
            _PRED_GREEN   = RGBColor(0xC6, 0xEF, 0xCE)
            _PRED_HEADER_NAMES = {
                'Прямая соединительная линия 10',
                'Прямоугольник 11',
                'Прямоугольник 12',
                'Рисунок 6',
            }

            # Пагинация: до 12 «логических» строк на слайд (data-row или
            # section-header — каждая считается за 1). Не оставляем section-
            # заголовок без последующих строк (orphan) — если он попал в самый
            # конец чанка, переносим в начало следующего.
            # Правка пользователя 2026-05-21: было 14, снижено до 12 — на
            # Бугры-3 страница 18/19 не влезала (длинный текст «Перечень
            # нарушений» переносился word_wrap'ом за пределы row.height).
            _PRED_ROWS_PER_PAGE = 12
            _pred_chunks = []
            _curr = []
            for _it in _pred_rows:
                _curr.append(_it)
                if len(_curr) >= _PRED_ROWS_PER_PAGE:
                    if isinstance(_curr[-1], dict):
                        _last = _curr.pop()
                        _pred_chunks.append(_curr)
                        _curr = [_last]
                    else:
                        _pred_chunks.append(_curr)
                        _curr = []
            if _curr:
                _pred_chunks.append(_curr)
            _pred_total_pages = len(_pred_chunks)

            def _pred_set_cell(cell, text, *, bold=False, size=9, align=PP_ALIGN.CENTER,
                               fill=None, color=None):
                if color is None:
                    color = _PRED_WHITE if (fill is not None and fill == _PRED_NAVY) else _PRED_BLACK
                cell.text = ""
                p = cell.text_frame.paragraphs[0]; p.alignment = align
                r = p.add_run(); r.text = str(text)
                r.font.name = "Calibri"
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = color
                cell.text_frame.word_wrap = True
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_top = Emu(20000); cell.margin_bottom = Emu(20000)
                cell.margin_left = Emu(30000); cell.margin_right = Emu(30000)
                if fill is not None:
                    cell.fill.solid(); cell.fill.fore_color.rgb = fill

            for _pg_idx, _chunk in enumerate(_pred_chunks):
                _pred_slide = clone_slide(prs, slide3)
                for _sh in list(_pred_slide.shapes):
                    if getattr(_sh, 'name', '') not in _PRED_HEADER_NAMES:
                        remove_shape(_sh)
                _title = f"ПРЕДПИСАНИЯ СК ({PROJECT_SHORT})"
                # 20.05.2026: суффикс «— N/M» убран по запросу пользователя.
                # Номер страницы N/M уже есть в правом верхнем углу шапки.
                for _sh in _pred_slide.shapes:
                    if getattr(_sh, 'name', '') == 'Прямоугольник 11':
                        set_text(_sh, _title); break

                _sw, _sh_h = prs.slide_width, prs.slide_height
                _ml = Emu(341998)
                _tbl_w = _sw - 2 * _ml
                _tbl_top = Emu(900000)
                _hdr_h = Emu(420000)
                # Адаптивная высота строки. Section-row (dict) занимает 1/3
                # от data-row, поэтому «эффективная» длина чанка для расчёта
                # _data_h = data_count + section_count/3. Без этого учёта
                # _data_h занижался: при 14 строках (включая sections) каждая
                # data-row получала avail_h/14, а реальная высота 14 строк
                # из data+section по факту = 14_data + 0.33_section *_data_h
                # < avail_h — место оставалось, но word_wrap длинных текстов
                # «Перечня нарушений» торчал за пределы. Учёт фактических
                # весов даёт нам _data_h, под который текст помещается.
                _avail_h = _sh_h - _tbl_top - Emu(300000) - _hdr_h
                _data_count = sum(1 for _x in _chunk if not isinstance(_x, dict))
                _section_count = sum(1 for _x in _chunk if isinstance(_x, dict))
                _eff_rows = _data_count + _section_count / 3.0
                _data_h = min(Emu(560000), int(_avail_h / max(1, _eff_rows)))
                _n_rows = 1 + len(_chunk)
                _tbl_h = _hdr_h + int(_data_h * _eff_rows)

                _gframe = _pred_slide.shapes.add_table(_n_rows, 8, _ml, _tbl_top, _tbl_w, _tbl_h)
                _ptbl = _gframe.table
                _ptbl.rows[0].height = _hdr_h
                # Section-row («Предписания строительного контроля ...») — в 3
                # раза ниже data-row (узкий разделитель), data-rows — обычные.
                _section_h = max(Emu(80000), int(_data_h / 3))
                for _ri in range(1, _n_rows):
                    _is_section = isinstance(_chunk[_ri - 1], dict)
                    _ptbl.rows[_ri].height = _section_h if _is_section else _data_h
                for _ci, _fr in enumerate(_PRED_FRACS):
                    _ptbl.columns[_ci].width = Emu(int(_tbl_w * _fr))

                # Шапка
                for _ci, _h in enumerate(_PRED_HDR):
                    _pred_set_cell(_ptbl.cell(0, _ci), _h, bold=True, size=10,
                                   align=PP_ALIGN.CENTER, fill=_PRED_NAVY)

                # Строки данных. Правки 19.05/20.05.2026:
                # — шрифт везде Calibri 9pt;
                # — «Срыв (дней)» (ci=5) bold + красный (критическая);
                # — «Отметка об устранении» (ci=6) — если содержит «Устранено»
                #   (дословно), заливка светло-зелёная;
                # — «Статус снятия» (ci=7) = «Не снято» → заливка светло-розовая,
                #   чёрный шрифт, не bold;
                # — section-row (dict {"section": "..."}) — merge 8 ячеек,
                #   navy-фон + белый bold Pt(11), заголовок раздела предписаний.
                for _ri, _row in enumerate(_chunk, start=1):
                    if isinstance(_row, dict):
                        _c0 = _ptbl.cell(_ri, 0)
                        _c7 = _ptbl.cell(_ri, 7)
                        _c0.merge(_c7)
                        _pred_set_cell(_c0, _row["section"], bold=True, size=11,
                                       align=PP_ALIGN.CENTER, fill=_PRED_NAVY)
                        continue
                    for _ci, _val in enumerate(_row):
                        _val_s = str(_val).strip()
                        _is_lag = (_ci == 5)
                        _is_otmetka_done = (_ci == 6 and "Устранено" in _val_s)
                        _is_snyatie_ne = (_ci == 7 and _val_s == "Не снято")
                        _color = _PRED_RED if _is_lag else None
                        _bold  = _is_lag
                        _fill  = None
                        if _is_otmetka_done:
                            _fill = _PRED_GREEN
                        elif _is_snyatie_ne:
                            _fill = _PRED_PINK
                        _pred_set_cell(
                            _ptbl.cell(_ri, _ci), _val,
                            bold=_bold,
                            size=9,
                            align=_PRED_ALIGNS[_ci],
                            fill=_fill,
                            color=_color,
                        )

                reorder_slide_to(prs, _pred_slide, len(prs.slides) - 1)
            print(f"Предписания slides added: {len(_pred_rows)} строк → {_pred_total_pages} слайд(а)")
        else:
            print("Предписания: нет просроченных позиций — слайд не добавляется")
    except Exception as _e:
        print(f"WARN: Предписания slide skipped: {_e}")
else:
    print("INFO: predpisaniya_summary_file_id не задан — слайд предписаний пропущен")

# Документы подрядчиков выносятся в отдельный портретный PDF и склеиваются
# с основным после convert_pdf.py (см. append_docs_portrait.py). В альбомном
# PPTX этот слайд не строится — A4 portrait нельзя смешать в одной презентации.

# Drop slides queued for removal (deferred until after clone_slide ops finish,
# so partname allocation in clone_slide doesn't collide with freed slot names)
_flush_slide_removals(prs)

# ---------------------------- renumber slides ----------------------------
import re as _re
total_n = len(prs.slides) + max(0, args.extra_pages)
slide_num_re = _re.compile(r'^\s*\d+\s*/\s*\d+\s*$')
for idx, sl in enumerate(prs.slides, start=1):
    for sh in sl.shapes:
        if not sh.has_text_frame: continue
        if slide_num_re.match(sh.text_frame.text.strip()):
            set_text(sh, f"{idx}/{total_n}")
            break

# Remove КОМПАКТ logo ('Рисунок 6') from every slide except the cover
for _sl in list(prs.slides)[1:]:
    for _sh in list(_sl.shapes):
        if getattr(_sh, 'name', '') == 'Рисунок 6':
            remove_shape(_sh)

# Убрать ` (PROJECT_SHORT)` суффикс из заголовков всех слайдов начиная со 2-го
# (на обложке отдельный шейп с полным названием — не трогаем).
# Здесь же — правка пользователя 20.05.2026:
#  — шрифт в шапке (Прямоугольник 11) увеличиваем ~ в 2 раза (Pt(36));
#  — выравниваем по центру.
# Шапка-шейп растёт вместе с текстом (a:spAutoFit), поэтому переразмещаем её
# (новый top/height) и сдвигаем всё ниже расположенное (separator + таблицы +
# скриншоты) на dy, чтобы избежать наезда на контент. Слайды-разделители
# (СМР/РД/...) не имеют Прямоугольник 11 (все шейпы сняты до подложки) —
# здесь они автоматически пропускаются.
# 20.05.2026: убрали якорь `$` — раньше на continuation-слайдах (СМР, РД, …
# когда стадия не помещается на одну страницу) заголовок был
# `"СТРОИТЕЛЬНО МОНТАЖНЫЕ РАБОТЫ ({PROJECT_SHORT}) — продолжение"`, и привязка
# к концу строки не давала вырезать `(PROJECT_SHORT)` в середине. Без `$`
# регэксп ловит суффикс где бы он ни стоял; PROJECT_SHORT в обычном тексте
# не встречается, ложных срабатываний нет.
_proj_suffix_re = _re.compile(r'\s*\(\s*' + _re.escape(PROJECT_SHORT) + r'\s*\)\s*')
_ORIG_TITLE_TOP    = 337182
_ORIG_TITLE_H      = 369332
_ORIG_TITLE_BOTTOM = _ORIG_TITLE_TOP + _ORIG_TITLE_H  # 706514
_NEW_TITLE_TOP     = 80000
_NEW_TITLE_H       = 720000
_NEW_TITLE_BOTTOM  = _NEW_TITLE_TOP + _NEW_TITLE_H    # 800000
_TITLE_DY          = _NEW_TITLE_BOTTOM - _ORIG_TITLE_BOTTOM  # ~93486
for _sl in list(prs.slides)[1:]:
    _title_sh = None
    for _sh in _sl.shapes:
        if getattr(_sh, 'name', '') == 'Прямоугольник 11':
            _title_sh = _sh; break
    if _title_sh is None or not _title_sh.has_text_frame: continue

    # 1) Текст: убрать суффикс «(PROJECT_SHORT)» через первый run (стили рана сохраняются),
    #    затем форсировать size=Pt(36), bold=True, центр.
    _full = _title_sh.text_frame.text
    _stripped = _proj_suffix_re.sub('', _full)
    for _p in _title_sh.text_frame.paragraphs:
        _p.alignment = PP_ALIGN.CENTER
        for _r in list(_p.runs):
            _r.text = ''
        if _p.runs and _stripped:
            _p.runs[0].text = _stripped
            _stripped = ''  # на остальные параграфы не пишем
        for _r in _p.runs:
            _r.font.size = Pt(36)
            _r.font.bold = True
            # 20.05.2026: явно проставляем navy. Раньше цвет наследовался из
            # theme (schemeClr tx2 + lumMod 75000) и в части стадийных
            # клонированных слайдов после set_text/strip-цикла мог сбрасываться
            # на чёрный. Жёстко фиксируем тот же navy, что у обложки.
            _r.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # 2) Геометрия: поднимаем верх шапки, удваиваем высоту. Сбрасываем spAutoFit,
    #    чтобы PowerPoint не пересчитывал размер контейнера.
    _title_sh.top    = Emu(_NEW_TITLE_TOP)
    _title_sh.height = Emu(_NEW_TITLE_H)
    _bodyPr = _title_sh.text_frame._txBody.find(qn('a:bodyPr'))
    if _bodyPr is not None:
        for _af in _bodyPr.findall(qn('a:spAutoFit')):
            _bodyPr.remove(_af)
        # Вертикальное центрирование текста внутри шейпа.
        _bodyPr.set('anchor', 'ctr')

    # 3) Слайд-номер (Прямоугольник 12) — синхронно с шапкой по Y.
    for _sh in _sl.shapes:
        if getattr(_sh, 'name', '') == 'Прямоугольник 12':
            _sh.top    = Emu(_NEW_TITLE_TOP)
            _sh.height = Emu(_NEW_TITLE_H)
            _bp2 = _sh.text_frame._txBody.find(qn('a:bodyPr'))
            if _bp2 is not None:
                _bp2.set('anchor', 'ctr')
            break

    # 4) Всё, что было ниже исходной шапки (separator, таблицы, скриншот ГПР)
    #    — сдвигаем на dy, чтобы сохранить исходный зазор шапка↔контент.
    for _sh in _sl.shapes:
        _n = getattr(_sh, 'name', '')
        if _n in ('Прямоугольник 11', 'Прямоугольник 12'): continue
        try:
            if _sh.top is not None and _sh.top >= _ORIG_TITLE_BOTTOM - 50000:
                _sh.top = _sh.top + _TITLE_DY
        except Exception:
            pass

os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT_PPTX)
print(f"Saved PPTX: {OUT_PPTX}  ({total_n} slides)")
print(f"Target PDF: {OUT_PDF}")
