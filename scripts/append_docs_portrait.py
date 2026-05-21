"""Дополняет основной PDF /мсг2 портретным слайдом «Документы подрядчиков».

Скачивает xlsx с Битрикс-Диска (id из cfg.docs_podryadchikov_file_id), делает
скриншот таблицы целиком через Excel COM, собирает A4 portrait PPTX, конвертит в
PDF через PowerPoint COM (ASCII temp из-за кириллических путей) и склеивает с
основным PDF: <original landscape pages> + <portrait page(s) at the end>.

Если строк больше PORTRAIT_ROWS_PER_PAGE — делит на несколько портретных слайдов,
шапка дублируется, со второго слайда в заголовке добавляется «(Продолжение)».

Запуск:
    python append_docs_portrait.py --config <cfg.json> --pdf <path-to-existing-PDF>
"""
import argparse, json, os, shutil, sys, tempfile, time

PORTRAIT_ROWS_PER_PAGE = 90   # сейчас 86 строк данных умещаются на одну страницу A4


def _read_webhook():
    path = r"C:\Авраменко\Claude Code Projects\.bitrix-webhook"
    with open(path, encoding="utf-8") as f:
        for line in f:
            s = line.strip()
            if not s or s.startswith("#"):
                continue
            wh = s.split()[0]
            break
    if not wh.startswith("http"):
        wh = f"https://kspb.bitrix24.ru/rest/2958/{wh}/"
    if not wh.endswith("/"):
        wh += "/"
    return wh


def bitrix_download_disk_file(file_id, out_path):
    import urllib.request, json as _json
    webhook = _read_webhook()
    info_url = f"{webhook}disk.file.get?id={file_id}"
    with urllib.request.urlopen(info_url, timeout=30) as r:
        data = _json.loads(r.read().decode("utf-8"))
    dl = data["result"]["DOWNLOAD_URL"]
    urllib.request.urlretrieve(dl, out_path)
    return out_path


def draw_outer_border(png_path, border_px=4, color=(0, 0, 0)):
    """Гарантированная рамка по краям PNG (Excel CopyPicture часто срезает
    выставленные xlEdge-границы)."""
    from PIL import Image, ImageDraw
    with Image.open(png_path) as im:
        im = im.convert("RGB")
        d = ImageDraw.Draw(im)
        w, h = im.size
        for i in range(border_px):
            d.rectangle([i, i, w - 1 - i, h - 1 - i], outline=color)
        im.save(png_path, "PNG")


def screenshot_multi(xlsx_path, sheet_name, ranges, out_paths, with_outer_border=True,
                     hidden_rows=None, hidden_cols=None):
    """Открыть xlsx в Excel COM, по очереди снять scrershot каждого range.
    with_outer_border=True — перед каждым CopyPicture обводит range внешней рамкой
    (xlEdgeLeft/Right/Top/Bottom = xlContinuous, xlMedium).

    Если переданы hidden_rows/hidden_cols — соответствующие строки и столбцы
    скрываются ДО скриншота через .Hidden=True. CopyPicture(Appearance=xlScreen)
    копирует только видимые ячейки, так что скрытое не попадёт в PNG."""
    import pythoncom, win32com.client
    from PIL import ImageGrab
    XL_EDGE_LEFT, XL_EDGE_TOP, XL_EDGE_BOTTOM, XL_EDGE_RIGHT = 7, 8, 9, 10
    XL_CONTINUOUS = 1
    XL_MEDIUM = -4138
    XL_AUTOMATIC = -4105
    pythoncom.CoInitialize()
    excel = win32com.client.DispatchEx("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False
    try:
        # ReadOnly=False — нам нужно ставить границы (рисование в БУ-памяти, не сохраняем)
        wb = excel.Workbooks.Open(xlsx_path, ReadOnly=False)
        ws = wb.Sheets(sheet_name)
        ws.Activate()
        try: excel.ActiveWindow.View = 1
        except Exception: pass
        # Скрываем заданные строки/колонки одним диапазоном на каждый отрезок
        # (быстрее, чем по одной).
        if hidden_rows:
            for s, e in _runs(sorted(hidden_rows)):
                ws.Range(f"{s}:{e}").EntireRow.Hidden = True
        if hidden_cols:
            from openpyxl.utils import get_column_letter as _gcl
            for s, e in _runs(sorted(hidden_cols)):
                ws.Range(f"{_gcl(s)}:{_gcl(e)}").EntireColumn.Hidden = True
        for addr, out_png in zip(ranges, out_paths):
            rng = ws.Range(addr)
            if with_outer_border:
                for edge in (XL_EDGE_LEFT, XL_EDGE_TOP, XL_EDGE_BOTTOM, XL_EDGE_RIGHT):
                    b = rng.Borders(edge)
                    b.LineStyle = XL_CONTINUOUS
                    b.Weight = XL_MEDIUM
                    b.ColorIndex = XL_AUTOMATIC
            rng.CopyPicture(Appearance=1, Format=2)
            time.sleep(0.6)
            img = ImageGrab.grabclipboard()
            if img is None:
                raise RuntimeError(f"Excel CopyPicture produced no clipboard image for {addr}")
            img.save(out_png, "PNG")
            if with_outer_border:
                draw_outer_border(out_png, border_px=4, color=(0, 0, 0))
        wb.Close(SaveChanges=False)
    finally:
        excel.Quit()
        pythoncom.CoUninitialize()


def _runs(sorted_ints):
    """[1,2,3,7,8,11] → [(1,3),(7,8),(11,11)]. Для группировки в Excel range."""
    if not sorted_ints:
        return []
    runs = []
    s = p = sorted_ints[0]
    for x in sorted_ints[1:]:
        if x == p + 1:
            p = x
        else:
            runs.append((s, p)); s = p = x
    runs.append((s, p))
    return runs


def compute_hidden_rows_cols(xlsx_path, data_top, data_bottom, left, right,
                              header_cols=2):
    """Вычисляет hidden_rows, hidden_cols для таблицы документов подрядчиков:
    скрываются те строки и столбцы, у которых нет ни одной отметки «нет» в
    data-области (правая часть от header_cols).

    header_cols — число «фиксированных» колонок-заголовков с левого края (A
    «№» и B «наименование документа»), они НЕ скрываются никогда.

    Шапка с подрядчиками (строки выше data_top) и левые колонки-заголовки
    остаются видимыми всегда."""
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    ws = wb.active
    def _is_net(v):
        if v is None: return False
        return str(v).strip().lower() == "нет"
    data_left = left + header_cols
    keep_cols = set(range(left, data_left))   # left header cols всегда видны
    for c in range(data_left, right + 1):
        for r in range(data_top, data_bottom + 1):
            if _is_net(ws.cell(row=r, column=c).value):
                keep_cols.add(c); break
    hidden_cols = [c for c in range(left, right + 1) if c not in keep_cols]
    hidden_rows = []
    for r in range(data_top, data_bottom + 1):
        has_net = any(_is_net(ws.cell(row=r, column=c).value)
                      for c in range(data_left, right + 1))
        if not has_net:
            hidden_rows.append(r)
    wb.close()
    return hidden_rows, hidden_cols


def detect_used_bbox(xlsx_path):
    """Возвращает (top_row, bottom_row, left_col, right_col) — реально непустой
    прямоугольник на первом листе (openpyxl, data_only). Колонки/строки, где ни
    в одной ячейке нет значения, отбрасываются с краёв."""
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    ws = wb.active
    rmax, cmax = ws.max_row, ws.max_column
    def col_has(c):
        return any(ws.cell(row=r, column=c).value not in (None, "") for r in range(1, rmax + 1))
    def row_has(r):
        return any(ws.cell(row=r, column=c).value not in (None, "") for c in range(1, cmax + 1))
    left = 1
    while left <= cmax and not col_has(left): left += 1
    right = cmax
    while right >= left and not col_has(right): right -= 1
    top = 1
    while top <= rmax and not row_has(top): top += 1
    bottom = rmax
    while bottom >= top and not row_has(bottom): bottom -= 1
    sheet_title = ws.title
    wb.close()
    return sheet_title, top, bottom, left, right


def convert_pptx_to_pdf_ascii(pptx_in, pdf_out):
    """PowerPoint COM SaveAs не работает с кириллическими путями — через ASCII temp."""
    import pythoncom, win32com.client
    ascii_dir = tempfile.gettempdir()
    ascii_pptx = os.path.join(ascii_dir, "_docs_portrait_tmp.pptx")
    ascii_pdf  = os.path.join(ascii_dir, "_docs_portrait_tmp.pdf")
    shutil.copyfile(pptx_in, ascii_pptx)
    for p in (ascii_pdf,):
        if os.path.exists(p):
            try: os.remove(p)
            except Exception: pass
    pythoncom.CoInitialize()
    pp = win32com.client.DispatchEx("PowerPoint.Application")
    try:
        d = pp.Presentations.Open(ascii_pptx, WithWindow=False)
        d.SaveAs(ascii_pdf, 32)  # ppSaveAsPDF
        d.Close()
    finally:
        pp.Quit()
        pythoncom.CoUninitialize()
    shutil.copyfile(ascii_pdf, pdf_out)


def merge_pdfs(base_pdf, append_pdf, out_pdf):
    from pypdf import PdfWriter, PdfReader
    w = PdfWriter()
    for src in (base_pdf, append_pdf):
        r = PdfReader(src)
        for page in r.pages:
            w.add_page(page)
    tmp = out_pdf + ".tmp"
    with open(tmp, "wb") as f:
        w.write(f)
    os.replace(tmp, out_pdf)


def build_portrait_pptx(png_pages, project_short, out_pptx,
                        revision_str=None, page_offset=0, total_pages=None):
    """Собрать A4-portrait PPTX с шапкой как у альбомных слайдов:
    - синий заголовок «ДОКУМЕНТЫ ПОДРЯДЧИКОВ» (без приписки проекта),
      на 2-м и далее слайдах — «ДОКУМЕНТЫ ПОДРЯДЧИКОВ (Продолжение)» (Продолжение в фуксии),
    - номер страницы справа «N/M» — если задан total_pages,
    - тонкая розовая горизонтальная линия (EC008C) под заголовком, без тени,
    - сам PNG таблицы по центру под линией.

    revision_str игнорируется (фидбэк 2026-05-18) — оставлен в сигнатуре."""
    from pptx import Presentation
    from pptx.util import Emu, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_CONNECTOR
    from lxml import etree
    from PIL import Image
    A4_W = Cm(21).emu
    A4_H = Cm(29.7).emu
    # 20.05.2026: подгоняем шапку портретных слайдов под альбомные — Pt(36)
    # bold Montserrat, navy `#1F4E79` (раньше тут был чуть более тёмный navy
    # `#1F3864`, что расходилось с альбомной обложкой). TITLE_TOP/TITLE_H
    # увеличены в ~2 раза, линия и таблица-картинка ниже сдвигаются автоматически.
    NAVY    = RGBColor(0x1F, 0x4E, 0x79)
    MAGENTA = RGBColor(0xEC, 0x00, 0x8C)
    TITLE_TOP = Emu(80000)
    TITLE_H   = Emu(720000)
    LINE_TOP  = Emu(860000)
    SIDE_PAD  = Emu(200000)
    LINE_LEFT  = SIDE_PAD
    LINE_WIDTH = A4_W - 2 * SIDE_PAD
    NUM_W = Emu(1287781)
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    prs = Presentation()
    prs.slide_width = A4_W
    prs.slide_height = A4_H
    blank = prs.slide_layouts[6]
    show_num = total_pages is not None
    tot = total_pages or 0

    def _style_run(run, *, text, color, size=14, bold=True, font="Montserrat"):
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color

    for idx, png in enumerate(png_pages, start=1):
        slide = prs.slides.add_slide(blank)

        # --- заголовок (центр)
        title_box = slide.shapes.add_textbox(Emu(0), TITLE_TOP, A4_W, TITLE_H)
        tf = title_box.text_frame
        tf.word_wrap = False
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        # Правка пользователя 2026-05-20: заголовок ÷1.4 (36 → 26pt) — раньше
        # был великоват для портретной A4-шапки, после уменьшения подравнивается
        # с подписями таблицы и оставляет больше воздуха.
        if idx == 1:
            _style_run(p.add_run(), text="ДОКУМЕНТЫ ПОДРЯДЧИКОВ", color=NAVY, size=26)
        else:
            _style_run(p.add_run(), text="ДОКУМЕНТЫ ПОДРЯДЧИКОВ ", color=NAVY, size=26)
            _style_run(p.add_run(), text="(Продолжение)", color=MAGENTA, size=26)

        # --- номер страницы N/M (правый край) — только если total_pages передан
        if show_num:
            num_box = slide.shapes.add_textbox(A4_W - NUM_W - SIDE_PAD, TITLE_TOP, NUM_W, TITLE_H)
            ntf = num_box.text_frame
            ntf.margin_top = ntf.margin_bottom = ntf.margin_right = 0
            ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
            np_ = ntf.paragraphs[0]
            np_.alignment = PP_ALIGN.RIGHT
            # Номер страницы остаётся компактным (Pt(14)) — это вспомогательный
            # элемент, его не масштабируем вместе с заголовком.
            _style_run(np_.add_run(), text=f"{page_offset + idx}/{tot}", color=NAVY, size=14)

        # --- розовая линия под заголовком (без тени)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          LINE_LEFT, LINE_TOP, LINE_LEFT + LINE_WIDTH, LINE_TOP)
        line.line.color.rgb = MAGENTA
        line.line.width = Emu(9525)  # 0.75pt — как в шаблоне
        # Убираем тень: (а) удаляем <p:style> — там effectRef из темы (с тенью),
        # (б) добавляем пустой <a:effectLst/> в <p:spPr> как явный «no effects».
        P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
        cxn = line._element
        for st in cxn.findall(f'{{{P_NS}}}style'):
            cxn.remove(st)
        sp_pr = cxn.find(f'{{{P_NS}}}spPr')
        if sp_pr is not None:
            for old in sp_pr.findall(f'{{{A_NS}}}effectLst'):
                sp_pr.remove(old)
            etree.SubElement(sp_pr, f'{{{A_NS}}}effectLst')

        # --- таблица-картинка
        with Image.open(png) as im:
            iw, ih = im.size
        ml = SIDE_PAD
        aw = A4_W - ml * 2
        top = LINE_TOP + Emu(200000)         # небольшой отступ от линии
        ah = A4_H - top - Emu(300000)
        scale = min(aw / iw, ah / ih)
        nw = int(iw * scale); nh = int(ih * scale)
        pl = (A4_W - nw) // 2
        pt = top + (ah - nh) // 2
        slide.shapes.add_picture(png, pl, pt, width=nw, height=nh)

    prs.save(out_pptx)


def _precount_portrait_pages(cfg):
    """Качает xlsx, считает число портретных страниц по тем же правилам, что
    и main(). Используется skill'ом /мсг2 для пре-вычисления --extra-pages
    в build_report.py до основной сборки."""
    file_id = cfg.get("docs_podryadchikov_file_id")
    if not file_id:
        return 0
    project_key = cfg.get("name", "проект").replace(" ", "_")
    tmp = tempfile.gettempdir()
    xlsx = os.path.join(tmp, f"docs_portrait_{project_key}.xlsx")
    bitrix_download_disk_file(file_id, xlsx)
    import openpyxl
    sheet, top, bottom, left, right = detect_used_bbox(xlsx)
    wb = openpyxl.load_workbook(xlsx, data_only=True)
    ws = wb.active
    top_non_empty = [(c, ws.cell(row=top, column=c).value)
                     for c in range(left, right + 1)
                     if ws.cell(row=top, column=c).value not in (None, "")]
    if len(top_non_empty) == 1 and "редакция" in str(top_non_empty[0][1]).lower():
        top += 1
    wb.close()
    HEADER_ROWS = 2
    data_start = top + HEADER_ROWS
    data_end = bottom
    # Правка пользователя 20.05.2026: строки/столбцы без отметок «нет»
    # скрываются. Число страниц считаем по ВИДИМЫМ data-строкам.
    hidden_rows, _ = compute_hidden_rows_cols(xlsx, data_start, data_end, left, right)
    n_visible = max(0, (data_end - data_start + 1) - len(hidden_rows))
    if n_visible == 0:
        return 1
    return (n_visible + PORTRAIT_ROWS_PER_PAGE - 1) // PORTRAIT_ROWS_PER_PAGE


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--config", required=True)
    ap.add_argument("--pdf", help="существующий итоговый PDF, к которому добавляется страница(ы)")
    ap.add_argument("--precount", action="store_true",
                    help="Вывести в stdout число портретных страниц (для pre-flight в /мсг2). PDF не нужен.")
    args = ap.parse_args()

    with open(args.config, encoding="utf-8") as f:
        cfg = json.load(f)

    if args.precount:
        print(_precount_portrait_pages(cfg))
        return

    file_id = cfg.get("docs_podryadchikov_file_id")
    if not file_id:
        print("INFO: docs_podryadchikov_file_id не задан — пропуск")
        return
    if not args.pdf:
        sys.exit("--pdf обязателен (или используйте --precount)")
    if not os.path.exists(args.pdf):
        sys.exit(f"PDF не найден: {args.pdf}")

    project_short = cfg.get("project_short", "")
    project_key = cfg.get("name", "проект").replace(" ", "_")

    tmp = tempfile.gettempdir()
    xlsx = os.path.join(tmp, f"docs_portrait_{project_key}.xlsx")
    bitrix_download_disk_file(file_id, xlsx)

    from openpyxl.utils import get_column_letter
    import openpyxl, re as _re
    sheet, top, bottom, left, right = detect_used_bbox(xlsx)

    # Caption-row (например R1 «редакция: 18.05.2026»): если в top-строке только
    # одна непустая ячейка с подстрокой "редакция" — снимаем её отдельно как
    # подзаголовок, в screenshot не включаем.
    revision_str = None
    wb = openpyxl.load_workbook(xlsx, data_only=True)
    ws = wb.active
    top_non_empty = [(c, ws.cell(row=top, column=c).value)
                     for c in range(left, right + 1)
                     if ws.cell(row=top, column=c).value not in (None, "")]
    if len(top_non_empty) == 1 and "редакция" in str(top_non_empty[0][1]).lower():
        m = _re.search(r"редакция[:\s]*(.+)", str(top_non_empty[0][1]), _re.I)
        revision_str = (m.group(1).strip() if m else str(top_non_empty[0][1]).strip())
        top += 1
    wb.close()

    col_first = get_column_letter(left)
    col_last  = get_column_letter(right)

    HEADER_ROWS = 2  # после среза caption-row: R2 (шапка с подрядчиками) + R3 (пусто)
    data_start = top + HEADER_ROWS
    data_end = bottom

    # Скрываем строки/столбцы без отметок «нет» (правка пользователя 20.05.2026).
    hidden_rows, hidden_cols = compute_hidden_rows_cols(
        xlsx, data_start, data_end, left, right)
    hidden_set = set(hidden_rows)
    visible_data_rows = [r for r in range(data_start, data_end + 1) if r not in hidden_set]
    n_visible = len(visible_data_rows)
    n_pages = max(1, (n_visible + PORTRAIT_ROWS_PER_PAGE - 1) // PORTRAIT_ROWS_PER_PAGE) \
              if n_visible else 1
    print(f"docs_portrait: всего data-строк={data_end - data_start + 1}, "
          f"скрыто строк={len(hidden_rows)}, скрыто колонок={len(hidden_cols)}, "
          f"видимо строк={n_visible}, страниц={n_pages}")

    r_header = f"{col_first}{top}:{col_last}{top + HEADER_ROWS - 1}"

    pages_png = []
    if n_pages == 1:
        # Одна страница — снимаем сразу всё одним range от шапки до последней
        # видимой data-строки. Скрытое (по hidden_rows/hidden_cols) Excel
        # пропустит при CopyPicture(Appearance=xlScreen).
        full_png = os.path.join(tmp, f"docs_portrait_full_{project_key}.png")
        end_row = visible_data_rows[-1] if visible_data_rows else data_end
        screenshot_multi(xlsx, sheet, [f"{col_first}{top}:{col_last}{end_row}"],
                         [full_png], hidden_rows=hidden_rows, hidden_cols=hidden_cols)
        pages_png.append(full_png)
    else:
        # Многостраничный режим — чанки по N видимых строк. Каждый chunk
        # адресуем диапазоном между ПЕРВОЙ и ПОСЛЕДНЕЙ видимой строкой
        # своего отрезка; скрытое внутри Excel выкинет.
        chunks = []
        for i in range(n_pages):
            s_idx = i * PORTRAIT_ROWS_PER_PAGE
            e_idx = min(s_idx + PORTRAIT_ROWS_PER_PAGE, n_visible) - 1
            s_row = visible_data_rows[s_idx]
            e_row = visible_data_rows[e_idx]
            chunks.append(f"{col_first}{s_row}:{col_last}{e_row}")
        png_header = os.path.join(tmp, f"docs_portrait_hdr_{project_key}.png")
        png_chunks = [os.path.join(tmp, f"docs_portrait_p{i+1}_{project_key}.png")
                      for i in range(n_pages)]
        screenshot_multi(xlsx, sheet, [r_header] + chunks, [png_header] + png_chunks,
                         hidden_rows=hidden_rows, hidden_cols=hidden_cols)
        from PIL import Image as _Img
        for i, pc in enumerate(png_chunks):
            out = os.path.join(tmp, f"docs_portrait_slide{i+1}_{project_key}.png")
            with _Img.open(png_header) as t, _Img.open(pc) as b:
                w = max(t.width, b.width)
                h = t.height + b.height
                im = _Img.new("RGB", (w, h), (255, 255, 255))
                im.paste(t, (0, 0))
                im.paste(b, (0, t.height))
                im.save(out, "PNG")
            draw_outer_border(out, border_px=4, color=(0, 0, 0))
            pages_png.append(out)

    # Для нумерации «N/M» на портретных слайдах нужно знать число альбомных страниц.
    # base_pages уже сидят в текущем PDF (build_report должен был учесть extra_pages
    # в renumber'е альбомных, чтобы M на всех слайдах совпадал).
    from pypdf import PdfReader as _PdfReader
    base_pages = len(_PdfReader(args.pdf).pages)
    total_pages = base_pages + len(pages_png)

    portrait_pptx = os.path.join(tmp, f"docs_portrait_{project_key}.pptx")
    portrait_pdf  = os.path.join(tmp, f"docs_portrait_{project_key}.pdf")
    build_portrait_pptx(pages_png, project_short, portrait_pptx,
                        page_offset=base_pages, total_pages=total_pages)
    convert_pptx_to_pdf_ascii(portrait_pptx, portrait_pdf)

    merge_pdfs(args.pdf, portrait_pdf, args.pdf)
    sys.stdout.buffer.write(f"Appended {n_pages} portrait page(s) -> {args.pdf}\n".encode("utf-8"))


if __name__ == "__main__":
    main()
