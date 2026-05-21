"""Build МСГ critical lag report for Бугры-3 (week 19) by analogy with Репино template."""
import sys
import os
import json
import shutil
import datetime as dt
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# Paths
TPL = r"C:\Users\amy\AppData\Local\Temp\repino_template.pptx"
COVER_IMG = r"C:\Users\amy\AppData\Local\Temp\bugri3_imgs\image1.png"
ITEMS_JSON = r"C:\Users\amy\AppData\Local\Temp\bugri3_items.json"
OUT_DIR = r"C:\Авраменко\1. КОМПАКТ\Отчеты\Еженедельные отчеты по пятницам\Бугры-3"
WEEK = 19
PERIOD = "04.05–10.05.2026"
TODAY = dt.date(2026, 5, 8)
PROJECT_NAME = "Складской комплекс «Бугры-3»"
PROJECT_ADDR = ("по адресу: Ленинградская обл., Всеволожский р-н, Бугровское г.п., "
                "земельный участок с кадастровым номером 47:07:0713003:16523")
PROJECT_SHORT = "СК «БУГРЫ-3»"
OUT_PPTX = os.path.join(OUT_DIR, f"Отчет из МСГ Бугры-3.pptx")
OUT_PDF  = os.path.join(OUT_DIR, f"2. МСГ критические отставания СК Бугры-3 неделя {WEEK}.pdf")

with open(ITEMS_JSON, encoding="utf-8") as f:
    ITEMS = json.load(f)

def _to_pct(v):
    if v is None: return None
    try:
        s = str(v).replace(',', '.').replace('+', '').replace('%', '').replace(' ', '').replace('\xa0', '').strip()
        return float(s)
    except: return None

# Compute lifecycle summary from data
def _avg_pct(rtypes):
    vals = [_to_pct(it['pct']) for it in ITEMS if it['rtype'] in rtypes]
    vals = [v for v in vals if v is not None]
    return sum(vals)/len(vals) if vals else 0.0

_stages = [
    (1, "Выпуск рабочей документации (РД)",          ['План РД'], ['Факт РД']),
    (2, "Формирование пакета документов (Пакет)",    ['План П'],  ['Факт П']),
    (3, "Тендер",                                     ['План Т'],  ['Факт Т']),
    (4, "Заключение договора",                        ['План Д'],  ['Факт Д']),
    (5, "Финансирование (оплата аванса)",             ['План Ф'],  ['Факт Ф']),
    (6, "Мобилизация подрядчика и МТР",               ['План М'],  ['Факт М']),
    (7, "Строительно-монтажные работы (СМР)",         ['План'],    ['Факт']),
]
LIFECYCLE = []
for n, name, pt, ft in _stages:
    p = _avg_pct(pt); f = _avg_pct(ft)
    LIFECYCLE.append((n, name, p, f, f - p))

print("Computed lifecycle:")
for n, nm, p, f, d in LIFECYCLE:
    print(f"  {n}. {nm}: План={p:.1f} Факт={f:.1f} Δ={d:+.1f}")

def parse_date(s):
    if not s: return None
    s = str(s).strip()
    if not s or s.startswith('#') or s == '-': return None
    for fmt in ("%Y-%m-%d", "%d.%m.%Y", "%d.%m.%y"):
        try:
            d = dt.datetime.strptime(s, fmt).date()
            if d.year < 2010: return None  # bogus epoch dates
            return d
        except: pass
    return None

# Pair adjacent rows: each "План X" is followed by its "Факт X"
def build_pairs():
    pairs = []
    i = 0
    while i < len(ITEMS):
        cur = ITEMS[i]
        if i+1 < len(ITEMS) and cur['rtype'].startswith('План') and ITEMS[i+1]['rtype'].startswith('Факт'):
            pairs.append({'plan': cur, 'fact': ITEMS[i+1]})
            i += 2
        else:
            i += 1
    return pairs

PAIRS = build_pairs()

# Filter helpers — collect lagging items per stage
def _to_float_safe(v):
    if v is None or v == '': return None
    try:
        if isinstance(v, str):
            s = v.replace('\xa0', '').replace(' ', '').replace(',', '.').replace('+', '').strip()
            return float(s)
        return float(v)
    except: return None

def is_completed(fact):
    """Both fact start and fact end are filled = work completed."""
    fn = parse_date(fact['pf_nach'])
    fo = parse_date(fact['pf_okon'])
    return fn is not None and fo is not None

def lagging_by_stage(plan_rtype):
    """Items with negative ОБЪЁМ ОТКЛ. on plan row, excluding completed."""
    out = []
    for p in PAIRS:
        if p['plan']['rtype'] != plan_rtype: continue
        plan = p['plan']; fact = p['fact']
        if is_completed(fact): continue
        n = _to_float_safe(plan.get('obj_otkl'))
        if n is None or n >= 0: continue
        plan_end = parse_date(plan['pf_okon'])
        fact_end = parse_date(fact['pf_okon'])
        days_dev = None
        if plan_end:
            days_dev = (plan_end - (fact_end or TODAY)).days
        out.append({
            'name': plan['name'], 'block': plan.get('block') or '',
            'section': plan.get('section') or '',
            'plan_n': plan['pf_nach'], 'plan_o': plan['pf_okon'],
            'fact_n': fact['pf_nach'], 'fact_o': fact['pf_okon'],
            'days_dev': days_dev,
            'obj_otkl': plan.get('obj_otkl'),
            'pct': fact.get('pct'),
        })
    return out

def lagging_smr():
    """СМР: items with negative ОБЪЁМ ОТКЛ., excluding completed."""
    out = []
    for p in PAIRS:
        if p['plan']['rtype'] != 'План': continue
        plan = p['plan']; fact = p['fact']
        if is_completed(fact): continue
        n = _to_float_safe(plan.get('obj_otkl'))
        if n is None or n >= 0: continue
        plan_end = parse_date(plan['pf_okon'])
        fact_end = parse_date(fact['pf_okon'])
        days_dev = None
        if plan_end:
            days_dev = (plan_end - (fact_end or TODAY)).days
        out.append({
            'name': plan['name'], 'block': plan.get('block') or '',
            'section': plan.get('section') or '',
            'plan_n': plan['pf_nach'], 'plan_o': plan['pf_okon'],
            'fact_n': fact['pf_nach'], 'fact_o': fact['pf_okon'],
            'days_dev': days_dev,
            'obj_otkl': plan.get('obj_otkl'),
            'unit': plan.get('unit'),
            'proj_obj': plan.get('proj_obj'),
        })
    return out

# Collect data
RD_LAG  = lagging_by_stage('План РД')
PAK_LAG = lagging_by_stage('План П')
TND_LAG = lagging_by_stage('План Т')
DOG_LAG = lagging_by_stage('План Д')
FIN_LAG = lagging_by_stage('План Ф')
MOB_LAG = lagging_by_stage('План М')
SMR_LAG = lagging_smr()

print(f"Lag counts: РД={len(RD_LAG)} Пакет={len(PAK_LAG)} Тендер={len(TND_LAG)} Договор={len(DOG_LAG)} Фин={len(FIN_LAG)} Моб={len(MOB_LAG)} СМР={len(SMR_LAG)}")

# Open template and modify
prs = Presentation(TPL)

# Helper: replace text in shape preserving first run formatting
def set_text(shape, new_text):
    if not shape.has_text_frame: return
    tf = shape.text_frame
    # Keep first paragraph, remove others
    p = tf.paragraphs[0]
    # Save first run formatting
    if p.runs:
        first_run = p.runs[0]
        # clear all runs
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        # Add new run with same formatting as first
        new_r = p.add_run()
        new_r.text = new_text
        # try to preserve font
        try:
            new_r.font.name = first_run.font.name
            new_r.font.size = first_run.font.size
            new_r.font.bold = first_run.font.bold
            if first_run.font.color and first_run.font.color.type is not None:
                new_r.font.color.rgb = first_run.font.color.rgb
        except: pass
    else:
        new_r = p.add_run()
        new_r.text = new_text
    # remove other paragraphs
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)

def find_shape_by_text(slide, fragment):
    for s in slide.shapes:
        if s.has_text_frame and fragment in s.text_frame.text:
            return s
    return None

def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)

def remove_pictures(slide):
    """Remove all picture shapes from slide."""
    to_remove = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    for s in to_remove:
        remove_shape(s)
    return to_remove

# ============ SLIDE 1: COVER ============
slide1 = prs.slides[0]
# Replace project name + address (Прямоугольник 5)
for shape in slide1.shapes:
    if shape.has_text_frame and 'Репин' in shape.text_frame.text:
        tf = shape.text_frame
        # Two paragraphs: name and address
        for i, p in enumerate(list(tf.paragraphs)):
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        # Set para 1
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = PROJECT_NAME
        r1.font.bold = True; r1.font.size = Pt(24); r1.font.name = 'Calibri'
        # Set para 2
        if len(tf.paragraphs) >= 2:
            p2 = tf.paragraphs[1]
        else:
            p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = PROJECT_ADDR
        r2.font.size = Pt(18); r2.font.name = 'Calibri'
        # Remove extra paragraphs
        for extra in tf.paragraphs[2:]:
            extra._p.getparent().remove(extra._p)
        break

# Replace banner with prominent two-line title — move higher, larger font
for shape in slide1.shapes:
    if shape.has_text_frame and 'НЕДЕЛЯ' in shape.text_frame.text:
        tf = shape.text_frame
        # Resize and move banner to occupy bigger area (center vertically, height ~1.6")
        shape.left = Emu(0)
        shape.top = Emu(2750000)
        shape.width = prs.slide_width
        shape.height = Emu(1500000)
        # Clear all existing paragraphs/runs
        for p in list(tf.paragraphs):
            p._p.getparent().remove(p._p)
        # Para 1: МСГ. КРИТИЧЕСКИЕ ОТСТАВАНИЯ
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = "МСГ. КРИТИЧЕСКИЕ ОТСТАВАНИЯ"
        r1.font.bold = True; r1.font.size = Pt(54); r1.font.name = 'Calibri'
        r1.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        # Para 2: НЕДЕЛЯ
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = f"НЕДЕЛЯ {WEEK}  ({PERIOD})"
        r2.font.bold = True; r2.font.size = Pt(40); r2.font.name = 'Calibri'
        r2.font.color.rgb = RGBColor(0xE6, 0x00, 0x73)
        break

# Replace cover image with Бугры-3 — full slide, 20% opacity, behind text
to_replace = None
for shape in slide1.shapes:
    if shape.shape_type == 13 and shape.name == 'Рисунок 4':
        to_replace = shape; break
if to_replace:
    remove_shape(to_replace)
# Insert background picture — full slide
bg_pic = slide1.shapes.add_picture(COVER_IMG, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
# Apply 80% transparency (so 20% visible) using alphaModFix
blip_fill = bg_pic._element.find('.//' + qn('p:blipFill'))
if blip_fill is not None:
    blip = blip_fill.find(qn('a:blip'))
    if blip is not None:
        alpha = etree.SubElement(blip, qn('a:alphaModFix'))
        alpha.set('amt', '20000')  # 20% opacity
# Move bg_pic to bottom of z-order
spTree = bg_pic._element.getparent()
spTree.remove(bg_pic._element)
insert_idx = 0
for i, child in enumerate(spTree):
    tag = etree.QName(child).localname
    if tag in ('nvGrpSpPr', 'grpSpPr'):
        insert_idx = i + 1
spTree.insert(insert_idx, bg_pic._element)

# ============ SLIDE 2: LIFECYCLE SUMMARY TABLE ============
slide2 = prs.slides[1]
# Update title
for s in slide2.shapes:
    if s.has_text_frame and 'СВОДНАЯ' in s.text_frame.text:
        set_text(s, f"СВОДНАЯ ТАБЛИЦА ЖИЗНЕННОГО ЦИКЛА ({PROJECT_SHORT})")
        break
# Remove main picture (Рисунок 4) — keep Рисунок 6 (logo)
for s in list(slide2.shapes):
    if s.shape_type == 13 and s.name == 'Рисунок 4':
        # save position
        L, T, W, H = s.left, s.top, s.width, s.height
        remove_shape(s)
        # Add native table at same position
        table_shape = slide2.shapes.add_table(rows=len(LIFECYCLE)+1, cols=5, left=L, top=T, width=W, height=H)
        tbl = table_shape.table
        # Column widths (proportional)
        total_w = W
        col_props = [0.06, 0.50, 0.14, 0.14, 0.16]
        for ci, prop in enumerate(col_props):
            tbl.columns[ci].width = int(total_w * prop)
        headers = ["№", "Этап жизненного цикла", "План, %", "Факт, %", "Δ (Факт−План), п.п."]
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = h
            r.font.bold = True; r.font.size = Pt(14); r.font.name = 'Calibri'
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for ri, row in enumerate(LIFECYCLE, start=1):
            num, name, plan, fact, delta = row
            vals = [str(num), name, f"{plan:.1f}".replace('.', ','), f"{fact:.1f}".replace('.', ','),
                    f"{'+' if delta>0 else ''}{delta:.1f}".replace('.', ',')]
            for ci, v in enumerate(vals):
                cell = tbl.cell(ri, ci)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                if ci == 1:
                    p.alignment = PP_ALIGN.LEFT
                else:
                    p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = v
                r.font.size = Pt(13); r.font.name = 'Calibri'
                if ci == 4:
                    r.font.bold = True
                    if delta < 0:
                        r.font.color.rgb = RGBColor(0xE0, 0x00, 0x00)
                    else:
                        r.font.color.rgb = RGBColor(0x00, 0x80, 0x00)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# ============ SLIDES 3-8: Update headings to use «БУГРЫ-3» ============
def update_title(slide, new_title):
    for s in slide.shapes:
        if s.has_text_frame and ('«РЕПИН»' in s.text_frame.text or 'РЕПИН' in s.text_frame.text or 'СКК' in s.text_frame.text):
            set_text(s, new_title)
            return

# Slide cloning + pagination helpers
def clone_slide(prs, src_slide, insert_after_idx=None):
    """Duplicate src_slide. Adds to end; caller may move via reorder_slides."""
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # Remove placeholder shapes added by layout
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)
    # Deep-copy shapes from source
    spTree = new_slide.shapes._spTree
    for shape in src_slide.shapes:
        new_el = deepcopy(shape._element)
        spTree.append(new_el)
    return new_slide

def reorder_slide_to(prs, slide, target_idx):
    """Move slide to target_idx in presentation."""
    sldIdLst = prs.slides._sldIdLst
    target_rId = None
    for rId, rel in prs.part.rels.items():
        if rel.target_part is slide.part:
            target_rId = rId; break
    for el in list(sldIdLst):
        if el.get(qn('r:id')) == target_rId:
            sldIdLst.remove(el)
            sldIdLst.insert(target_idx, el)
            return

def paginate_by_section(items, max_rows):
    """Greedy-pack items into pages, respecting section grouping. Returns list[list[item]]."""
    if not items: return []
    pages = []
    cur = []
    cur_rows = 1  # header
    last_section = None
    for it in items:
        sec = it.get('section') or '— БЕЗ РАЗДЕЛА —'
        new_section = sec != last_section
        rows_needed = (1 if new_section else 0) + 1
        if cur_rows + rows_needed > max_rows and cur:
            pages.append(cur)
            cur = []
            cur_rows = 1
            last_section = None
            new_section = True
            rows_needed = 2
        if new_section:
            last_section = sec
        cur.append(it)
        cur_rows += rows_needed
    if cur:
        pages.append(cur)
    return pages

def build_paged_table(prs, base_slide, items, columns, top_emu, base_title, title_setter, max_rows=24, subtitle_fn=None):
    """Render items across one or more slides. Returns list of slides used."""
    pages = paginate_by_section(items, max_rows=max_rows)
    if not pages:
        # Empty — just use base_slide
        build_data_table(base_slide, [], columns, top_emu=top_emu)
        return [base_slide]
    slides = []
    base_idx = list(prs.slides).index(base_slide)
    for pi, page_items in enumerate(pages):
        if pi == 0:
            target = base_slide
        else:
            target = clone_slide(prs, base_slide)
            # Move target right after previous slide
            reorder_slide_to(prs, target, base_idx + pi)
            # Update title with "(продолжение)"
            title_setter(target, base_title + " — продолжение")
        if subtitle_fn:
            subtitle_fn(target)
        build_data_table(target, page_items, columns, top_emu=top_emu)
        slides.append(target)
    return slides

# Universal data table builder with section grouping
def build_data_table(slide, items, columns, top_emu=800000):
    """columns is list of (header, key, width_prop, fmt). Groups items by 'section' field."""
    # remove pictures (except logo) AND any existing tables (from cloned slide)
    for s in list(slide.shapes):
        if s.shape_type == 13 and s.name != 'Рисунок 6':  # keep logo
            remove_shape(s)
        elif s.has_table:
            remove_shape(s)
    # position: full width below title
    sw = prs.slide_width
    sh = prs.slide_height
    L = Emu(341998)
    T = Emu(top_emu)
    W = Emu(sw - L*2)
    Hpx = sh - T - Emu(300000)
    H = Hpx
    if not items:
        # No items
        tb = slide.shapes.add_textbox(L, Emu(2942274), W, Emu(369332))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "ОТСТАВАНИЯ ОТ ПЛАНА ОТСУТСТВУЮТ"
        r.font.bold = True; r.font.size = Pt(24); r.font.name = 'Calibri'
        r.font.color.rgb = RGBColor(0xE6, 0x00, 0x73)
        return
    # Group by section preserving order
    grouped = []  # list of (section, [items])
    for it in items:
        sec = it.get('section') or '— БЕЗ РАЗДЕЛА —'
        if grouped and grouped[-1][0] == sec:
            grouped[-1][1].append(it)
        else:
            grouped.append((sec, [it]))
    # Compute total table rows: 1 header + sum(section_header + items per group)
    total_rows = 1 + sum(1 + len(grp) for _, grp in grouped)
    table_shape = slide.shapes.add_table(rows=total_rows, cols=len(columns), left=L, top=T, width=W, height=H)
    tbl = table_shape.table
    total_w = W
    props = [c[2] for c in columns]
    sp = sum(props); props = [p/sp for p in props]
    for ci, p in enumerate(props):
        tbl.columns[ci].width = int(total_w * p)
    # Header
    for ci, (hdr, key, _, fmt) in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(10); r.font.name = 'Calibri'
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Body with section grouping
    ri = 0
    for sec, grp in grouped:
        ri += 1
        # Section header row — merge across all columns
        first = tbl.cell(ri, 0); last = tbl.cell(ri, len(columns)-1)
        first.merge(last)
        first.text = ""
        p = first.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = sec
        r.font.bold = True; r.font.size = Pt(10); r.font.name = 'Calibri'
        r.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        first.fill.solid()
        first.fill.fore_color.rgb = RGBColor(0xDC, 0xE6, 0xF1)
        first.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Items in section
        for item in grp:
            ri += 1
            for ci, (hdr, key, _, fmt) in enumerate(columns):
                cell = tbl.cell(ri, ci)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                v = item.get(key)
                txt = fmt(v) if fmt else (str(v) if v not in (None,'') else '–')
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                rn = p.add_run(); rn.text = txt
                rn.font.size = Pt(9); rn.font.name = 'Calibri'
                if key == 'days_dev' and v is not None and isinstance(v, (int, float)) and v < 0:
                    rn.font.bold = True
                    rn.font.color.rgb = RGBColor(0xE0, 0x00, 0x00)
                if key == 'obj_otkl':
                    nv = _to_float_safe(v)
                    if nv is not None and nv < 0:
                        rn.font.bold = True
                        rn.font.color.rgb = RGBColor(0xE0, 0x00, 0x00)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def fmt_date(s):
    if not s: return '–'
    if isinstance(s, str) and (s.startswith('#') or s == '-'):
        return '–'
    d = parse_date(s)
    if d:
        return d.strftime("%d.%m.%y")
    return '–'

def fmt_int(v):
    if v is None or v == '': return '–'
    try:
        n = int(v)
        if n == 0: return '0'
        return f"{n:+d}"
    except: return str(v)

def _to_float(v):
    if isinstance(v, str):
        s = v.replace('\xa0', '').replace(' ', '').replace(',', '.').replace('+', '').strip()
        return float(s)
    return float(v)

def fmt_num(v):
    """For deviation columns: signed."""
    if v is None or v == '': return '–'
    try:
        f = _to_float(v)
        if abs(f) < 0.005: return '0'
        if abs(f - round(f)) < 1e-6:
            n = int(round(f))
            return f"{n:+d}"
        sign = '+' if f > 0 else ''
        return f"{sign}{f:.1f}".replace('.', ',')
    except: return str(v)

def fmt_pos(v):
    """For volume columns: no sign."""
    if v is None or v == '': return '–'
    try:
        f = _to_float(v)
        if abs(f - round(f)) < 1e-6:
            return str(int(round(f)))
        return f"{f:.1f}".replace('.', ',')
    except: return str(v)

# Pre-capture slide references BEFORE any pagination shifts indices
slide3 = prs.slides[2]   # СМР
slide4 = prs.slides[3]   # РД
slide5 = prs.slides[4]   # Тендер
slide6 = prs.slides[5]   # Договор
slide7 = prs.slides[6]   # Финансирование
slide8 = prs.slides[7]   # Мобилизация

# ---- SLIDE 3: СМР ----
update_title(slide3, f"СТРОИТЕЛЬНО МОНТАЖНЫЕ РАБОТЫ ({PROJECT_SHORT})")

def _set_title_smr(slide, text):
    for s in slide.shapes:
        if s.has_text_frame and ('СТРОИТЕЛЬНО' in s.text_frame.text.upper() or 'МОНТАЖНЫЕ' in s.text_frame.text.upper()):
            set_text(s, text); return

build_paged_table(prs, slide3, SMR_LAG, [
    ("Наименование работ", "name", 0.30, lambda v: str(v)[:80]),
    ("Здан.", "block", 0.05, lambda v: str(v) if v else '–'),
    ("Объём\nоткл.", "obj_otkl", 0.08, fmt_num),
    ("Ед.", "unit", 0.04, lambda v: str(v) if v else '–'),
    ("Проект.\nобъём", "proj_obj", 0.08, fmt_pos),
    ("План\nначало", "plan_n", 0.09, fmt_date),
    ("План\nоконч.", "plan_o", 0.09, fmt_date),
    ("Факт\nначало", "fact_n", 0.09, fmt_date),
    ("Факт\nоконч.", "fact_o", 0.09, fmt_date),
    ("Дни\nоткл.", "days_dev", 0.07, fmt_int),
], top_emu=800000, base_title=f"СТРОИТЕЛЬНО МОНТАЖНЫЕ РАБОТЫ ({PROJECT_SHORT})", title_setter=_set_title_smr, max_rows=24)

# ---- SLIDE 4: РД ----
update_title(slide4, f"РАБОЧАЯ ДОКУМЕНТАЦИЯ ({PROJECT_SHORT})")
# Replace text shape
for s in list(slide4.shapes):
    if s.has_text_frame and 'ПОЛНОМ' in s.text_frame.text:
        remove_shape(s)
# Add subtitle below title bar
sw = prs.slide_width
tb = slide4.shapes.add_textbox(Emu(0), Emu(770000), sw, Emu(420000))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "ВЫПУСК РД ОТСТАЁТ ОТ ПЛАНА НА 12,5 п.п. (ПЛАН 50,8% / ФАКТ 38,3%)"
r.font.bold = True; r.font.size = Pt(18); r.font.name = 'Calibri'
r.font.color.rgb = RGBColor(0xE6, 0x00, 0x73)
def _rd_subtitle(slide):
    sw = prs.slide_width
    tb = slide.shapes.add_textbox(Emu(0), Emu(770000), sw, Emu(420000))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "ВЫПУСК РД ОТСТАЁТ ОТ ПЛАНА НА 12,5 п.п. (ПЛАН 50,8% / ФАКТ 38,3%)"
    r.font.bold = True; r.font.size = Pt(18); r.font.name = 'Calibri'
    r.font.color.rgb = RGBColor(0xE6, 0x00, 0x73)
# Remove the subtitle we already added before pagination (we add fresh per page)
for s in list(slide4.shapes):
    if s.has_text_frame and 'ОТСТАЁТ' in s.text_frame.text:
        remove_shape(s)

def _set_title(slide, text):
    for s in slide.shapes:
        if s.has_text_frame and ('РАБОЧАЯ' in s.text_frame.text.upper() or 'РАБОЧАЙ' in s.text_frame.text.upper() or 'БУГРЫ' in s.text_frame.text):
            set_text(s, text)
            return

build_paged_table(prs, slide4, RD_LAG, [
    ("Наименование работ", "name", 0.36, lambda v: str(v)[:100]),
    ("Объём\nоткл.", "obj_otkl", 0.10, fmt_num),
    ("План\nначало", "plan_n", 0.10, fmt_date),
    ("План\nоконч.", "plan_o", 0.10, fmt_date),
    ("Факт\nначало", "fact_n", 0.10, fmt_date),
    ("Факт\nоконч.", "fact_o", 0.10, fmt_date),
    ("Дни\nоткл.", "days_dev", 0.08, fmt_int),
], top_emu=1250000, base_title=f"РАБОЧАЯ ДОКУМЕНТАЦИЯ ({PROJECT_SHORT})", title_setter=_set_title, max_rows=22, subtitle_fn=_rd_subtitle)

# ---- SLIDE 5: ТЕНДЕР ----
update_title(slide5, f"ТЕНДЕР ({PROJECT_SHORT})")
# Update status line text
for s in slide5.shapes:
    if s.has_text_frame and 'ОТСТАВАНИЯ' in s.text_frame.text:
        set_text(s, "ОТСТАВАНИЯ ОТ ПЛАНА ОТСУТСТВУЮТ")

# ---- SLIDE 6: ДОГОВОР ----
update_title(slide6, f"ДОГОВОР ({PROJECT_SHORT})")
# Remove vertical line on slide 6 (Прямая соединительная линия 4) before clone
for s in list(slide6.shapes):
    if hasattr(s, 'name') and 'Прямая соединительная линия 4' == s.name:
        remove_shape(s)

def _set_title_dog(slide, text):
    for s in slide.shapes:
        if s.has_text_frame and ('ДОГОВОР' in s.text_frame.text.upper()):
            set_text(s, text); return

build_paged_table(prs, slide6, DOG_LAG, [
    ("Наименование работ", "name", 0.36, lambda v: str(v)[:100]),
    ("Объём\nоткл.", "obj_otkl", 0.10, fmt_num),
    ("План\nначало", "plan_n", 0.10, fmt_date),
    ("План\nоконч.", "plan_o", 0.10, fmt_date),
    ("Факт\nначало", "fact_n", 0.10, fmt_date),
    ("Факт\nоконч.", "fact_o", 0.10, fmt_date),
    ("Дни\nоткл.", "days_dev", 0.08, fmt_int),
], top_emu=800000, base_title=f"ДОГОВОР ({PROJECT_SHORT})", title_setter=_set_title_dog, max_rows=24)

# ---- SLIDE 7: ФИНАНСИРОВАНИЕ ----
update_title(slide7, f"ФИНАНСИРОВАНИЕ ({PROJECT_SHORT})")
# Remove placeholder text
for s in list(slide7.shapes):
    if s.has_text_frame and ('ИНФОРМАЦИЯ' in s.text_frame.text or 'НЕ ЗАПОЛН' in s.text_frame.text or 'АВАНСЫ' in s.text_frame.text):
        remove_shape(s)

def _set_title_fin(slide, text):
    for s in slide.shapes:
        if s.has_text_frame and ('ФИНАНСИРОВАНИЕ' in s.text_frame.text.upper()):
            set_text(s, text); return

build_paged_table(prs, slide7, FIN_LAG, [
    ("Наименование работ", "name", 0.36, lambda v: str(v)[:100]),
    ("Объём\nоткл.", "obj_otkl", 0.10, fmt_num),
    ("План\nначало", "plan_n", 0.10, fmt_date),
    ("План\nоконч.", "plan_o", 0.10, fmt_date),
    ("Факт\nначало", "fact_n", 0.10, fmt_date),
    ("Факт\nоконч.", "fact_o", 0.10, fmt_date),
    ("Дни\nоткл.", "days_dev", 0.08, fmt_int),
], top_emu=800000, base_title=f"ФИНАНСИРОВАНИЕ ({PROJECT_SHORT})", title_setter=_set_title_fin, max_rows=24)

# ---- SLIDE 8: МОБИЛИЗАЦИЯ ----
update_title(slide8, f"МОБИЛИЗАЦИЯ ({PROJECT_SHORT})")
for s in list(slide8.shapes):
    if hasattr(s, 'name') and 'Прямая соединительная линия 4' == s.name:
        remove_shape(s)

def _set_title_mob(slide, text):
    for s in slide.shapes:
        if s.has_text_frame and ('МОБИЛИЗАЦИЯ' in s.text_frame.text.upper()):
            set_text(s, text); return

build_paged_table(prs, slide8, MOB_LAG, [
    ("Наименование работ", "name", 0.36, lambda v: str(v)[:100]),
    ("Объём\nоткл.", "obj_otkl", 0.10, fmt_num),
    ("План\nначало", "plan_n", 0.10, fmt_date),
    ("План\nоконч.", "plan_o", 0.10, fmt_date),
    ("Факт\nначало", "fact_n", 0.10, fmt_date),
    ("Факт\nоконч.", "fact_o", 0.10, fmt_date),
    ("Дни\nоткл.", "days_dev", 0.08, fmt_int),
], top_emu=800000, base_title=f"МОБИЛИЗАЦИЯ ({PROJECT_SHORT})", title_setter=_set_title_mob, max_rows=24)

# Renumber slides — find shapes containing "X/Y" pattern (usually in "Прямоугольник 12")
import re as _re
total_n = len(prs.slides)
slide_num_re = _re.compile(r'^\s*\d+\s*/\s*\d+\s*$')
for idx, sl in enumerate(prs.slides, start=1):
    for sh in sl.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text.strip()
        if slide_num_re.match(txt):
            set_text(sh, f"{idx}/{total_n}")
            break

os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT_PPTX)
print(f"Saved PPTX: {OUT_PPTX}  ({total_n} slides)")
print(f"Will save PDF to: {OUT_PDF}")
